import os
from datetime import datetime, timezone, timedelta
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_test_pptx():
    prs = Presentation()
    prs.slide_width = Cm(21.0)
    prs.slide_height = Cm(29.7)
    blank_slide_layout = prs.slide_layouts[6]
    
    # --- Cover Page ---
    slide = prs.slides.add_slide(blank_slide_layout)
    # Just text for now
    txBox = slide.shapes.add_textbox(Cm(2), Cm(20), Cm(17), Cm(5))
    tf = txBox.text_frame
    tf.text = "LAPORAN WORK ORDER"
    
    prs.save("test.pptx")
    print("test.pptx created")

if __name__ == "__main__":
    create_test_pptx()
