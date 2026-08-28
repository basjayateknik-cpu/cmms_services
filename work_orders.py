import os
import json
from flask import Blueprint, render_template, request, redirect, url_for, flash, current_app, Response, send_file
from werkzeug.utils import secure_filename
from flask_login import login_required, current_user
from models import db, WorkOrder, Asset, User, Part, StockLevel, WorkOrderPart, WorkOrderStatus, WorkOrderLog, ProjectCode, Tasklist, WorkOrderProcedure, WorkOrderChecklistParameter, Checklist, ChecklistParameterTemplate, Team, Site
from datetime import datetime, timezone
from utils import supervisor_or_admin_required, send_whatsapp_notification
import pandas as pd
import io
import tempfile
from fpdf import FPDF
from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def check_wo_access(wo):
    if current_user.role == 'Technician' and current_user.id not in [a.id for a in wo.assignees]:
        return False
    return True

def get_all_project_codes():
    return ProjectCode.query.order_by(ProjectCode.code).all()

work_orders_bp = Blueprint('work_orders', __name__, url_prefix='/work-orders')

@work_orders_bp.route('/')
@login_required
def index():
    from models import Site, SavedWorkOrderView
    status_filter = request.args.get('status')
    assigned_filter = request.args.get('assigned_to')
    site_filter = request.args.get('site_id', type=int)
    mt_filter = request.args.get('maintenance_type')
    search_query = request.args.get('q', '').strip()
    completeness_filter = request.args.get('completeness')
    
    from sqlalchemy.orm import selectinload
    query = WorkOrder.query.options(
        selectinload(WorkOrder.asset).selectinload(Asset.site),
        selectinload(WorkOrder.team),
        selectinload(WorkOrder.current_status),
        selectinload(WorkOrder.tasklist)
    )
    
    # Time Range Filtering
    time_range = request.args.get('time_range', 'this_month')
    from datetime import datetime, timezone
    import calendar
    now = datetime.now(timezone.utc).replace(tzinfo=None)
    start_date = None
    end_date = None
    
    if time_range == 'this_month':
        start_date = datetime(now.year, now.month, 1)
        _, last_day = calendar.monthrange(now.year, now.month)
        end_date = datetime(now.year, now.month, last_day, 23, 59, 59)
    elif time_range == 'last_month':
        last_month = now.month - 1 if now.month > 1 else 12
        year = now.year if now.month > 1 else now.year - 1
        start_date = datetime(year, last_month, 1)
        _, last_day = calendar.monthrange(year, last_month)
        end_date = datetime(year, last_month, last_day, 23, 59, 59)
    elif time_range == 'this_year':
        start_date = datetime(now.year, 1, 1)
        end_date = datetime(now.year, 12, 31, 23, 59, 59)
    elif time_range == 'custom':
        start_str = request.args.get('start_date')
        end_str = request.args.get('end_date')
        if start_str and end_str:
            try:
                start_date = datetime.strptime(start_str, '%Y-%m-%d')
                end_date = datetime.strptime(end_str, '%Y-%m-%d').replace(hour=23, minute=59, second=59)
            except ValueError:
                pass
                
    if start_date and end_date:
        from sqlalchemy import func
        query = query.filter(func.coalesce(WorkOrder.suggested_completion_date, WorkOrder.date_created) >= start_date, func.coalesce(WorkOrder.suggested_completion_date, WorkOrder.date_created) <= end_date)
        
    # Site filtering
    if current_user.site_id:
        query = query.join(Asset).filter(Asset.site_id == current_user.site_id)
    elif site_filter:
        query = query.join(Asset).filter(Asset.site_id == site_filter)
    
    # RBAC: Technicians only see their own assigned work orders
    if current_user.role == 'Technician':
        query = query.filter(WorkOrder.assignees.any(id=current_user.id))
    
    if status_filter:
        if status_filter in ['Active', 'Pending', 'Closed', 'Draft']:
            query = query.join(WorkOrderStatus).filter(WorkOrderStatus.control_type == status_filter)
        else:
            if str(status_filter).isdigit():
                query = query.filter(WorkOrder.status_id == int(status_filter))
            else:
                query = query.join(WorkOrderStatus).filter(WorkOrderStatus.name == status_filter)
    if assigned_filter:
        query = query.filter(WorkOrder.assignees.any(id=assigned_filter))
    if mt_filter:
        query = query.filter(WorkOrder.maintenance_type == mt_filter)
    if completeness_filter == 'no_tasklist':
        query = query.filter(WorkOrder.tasklist_id == None)
    elif completeness_filter == 'no_checklist':
        query = query.filter(WorkOrder.checklist_id == None)
    elif completeness_filter == 'no_both':
        query = query.filter(db.and_(WorkOrder.tasklist_id == None, WorkOrder.checklist_id == None))
    if search_query:
        search_term = f"%{search_query}%"
        query = query.filter(
            db.or_(
                WorkOrder.code.ilike(search_term),
                WorkOrder.description.ilike(search_term)
            )
        )
        
    work_orders = query.all()
    statuses = WorkOrderStatus.query.all()
    users = User.query.all()
    
    # Get available sites for filter (only for users without fixed site)
    available_sites = []
    if not current_user.site_id:
        available_sites = Site.query.order_by(Site.name).all()
        
    # Get unique maintenance types
    db_types = db.session.query(WorkOrder.maintenance_type).distinct().all()
    maintenance_types = [dt[0] for dt in db_types if dt[0]]
    if not maintenance_types:
        maintenance_types = ['Corrective', 'Preventive']
    elif 'Corrective' not in maintenance_types:
        maintenance_types.append('Corrective')
    elif 'Preventive' not in maintenance_types:
        maintenance_types.append('Preventive')
    
    view_mode = request.args.get('view', 'list')
    
    # Prepare JSON data for advanced views
    work_orders_data = []
    for wo in work_orders:
        work_orders_data.append({
            'id': wo.id,
            'code': wo.code,
            'description': wo.description,
            'priority': wo.priority or '-',
            'maintenance_type': wo.maintenance_type or '-',
            'asset': wo.asset.name if wo.asset else '-',
            'site': wo.asset.site.name if wo.asset and wo.asset.site else '-',
            'assignees': [a.name for a in wo.assignees] if wo.assignees else [],
            'team': wo.team.name if wo.team else '-',
            'status': wo.current_status.name if wo.current_status else 'Draft',
            'status_control': wo.current_status.control_type if wo.current_status else 'Draft',
            'project': wo.project_code or '-',
            'tasklist': wo.tasklist.name if wo.tasklist else '-',
            'customer_name': wo.customer_name or '-',
            'estimated_hours': wo.estimated_hours or 0,
            'date_created': str(wo.date_created.date()) if wo.date_created else '-',
            'start_date': str(wo.start_date.date()) if wo.start_date else '-',
            'target_date': str(wo.suggested_completion_date.date()) if wo.suggested_completion_date else '-',
        })
    work_orders_json = json.dumps(work_orders_data)
    start_date_str = start_date.strftime('%Y-%m-%d') if start_date else ''
    end_date_str = end_date.strftime('%Y-%m-%d') if end_date else ''
    
    # Fetch saved views
    saved_views = SavedWorkOrderView.query.filter(
        db.or_(SavedWorkOrderView.is_public == True, SavedWorkOrderView.user_id == current_user.id)
    ).order_by(SavedWorkOrderView.name).all()

    return render_template('work_orders/index.html', work_orders=work_orders, statuses=statuses, users=users, available_sites=available_sites, selected_site_id=site_filter, maintenance_types=maintenance_types, selected_mt=mt_filter, search_query=search_query, completeness_filter=completeness_filter, view_mode=view_mode, work_orders_json=work_orders_json, time_range=time_range, start_date_str=start_date_str, end_date_str=end_date_str, saved_views=saved_views)

@work_orders_bp.route('/save-view', methods=['POST'])
@login_required
def save_view():
    from models import SavedWorkOrderView
    from flask import jsonify
    name = request.form.get('name')
    is_public = request.form.get('is_public') == 'true'
    view_mode = request.form.get('view_mode', 'pivot')
    pivot_row = request.form.get('pivot_row', '')
    pivot_col = request.form.get('pivot_col', '')
    filters = request.form.get('filters', '{}')
    
    if not name:
        return jsonify({'success': False, 'message': 'Name is required'}), 400
        
    new_view = SavedWorkOrderView(
        name=name,
        user_id=current_user.id,
        is_public=is_public,
        view_mode=view_mode,
        pivot_row=pivot_row,
        pivot_col=pivot_col,
        filters=filters
    )
    db.session.add(new_view)
    db.session.commit()
    
    return jsonify({'success': True, 'id': new_view.id, 'name': new_view.name})

@work_orders_bp.route('/delete-view/<int:id>', methods=['POST'])
@login_required
def delete_view(id):
    from models import SavedWorkOrderView
    from flask import jsonify
    view = SavedWorkOrderView.query.get_or_404(id)
    if view.user_id != current_user.id and current_user.role not in ['Admin', 'Supervisor']:
        return jsonify({'success': False, 'message': 'Unauthorized'}), 403
        
    db.session.delete(view)
    db.session.commit()
    return jsonify({'success': True})

@work_orders_bp.route('/export')
@login_required
@supervisor_or_admin_required
def export_all():
    from models import Site
    status_filter = request.args.get('status')
    assigned_filter = request.args.get('assigned_to')
    site_filter = request.args.get('site_id', type=int)
    mt_filter = request.args.get('maintenance_type')
    format = request.args.get('format', 'excel')
    
    query = WorkOrder.query
    
    # Site filtering
    if current_user.site_id:
        query = query.join(Asset).filter(Asset.site_id == current_user.site_id)
    elif site_filter:
        query = query.join(Asset).filter(Asset.site_id == site_filter)
    
    # RBAC: Technicians only see their own assigned work orders (though supervisor_or_admin_required is active)
    if current_user.role == 'Technician':
        query = query.filter(WorkOrder.assignees.any(id=current_user.id))
    
    if status_filter:
        if status_filter in ['Active', 'Pending', 'Closed', 'Draft']:
            query = query.join(WorkOrderStatus).filter(WorkOrderStatus.control_type == status_filter)
        else:
            if str(status_filter).isdigit():
                query = query.filter(WorkOrder.status_id == int(status_filter))
            else:
                query = query.join(WorkOrderStatus).filter(WorkOrderStatus.name == status_filter)
    if assigned_filter:
        query = query.filter(WorkOrder.assignees.any(id=assigned_filter))
    if mt_filter:
        query = query.filter(WorkOrder.maintenance_type == mt_filter)
        
    work_orders = query.all()
    
    data = []
    for wo in work_orders:
        # Detailed sub-data concatenation
        procedures_list = []
        for p in wo.procedures:
            status = "Done" if p.is_completed else "Pending"
            procedures_list.append(f"- {p.name} ({status})")
        procedures_text = "\n".join(procedures_list)

        checklist_list = []
        for c in wo.checklist_parameters:
            val = c.value if c.value else "N/A"
            checklist_list.append(f"- {c.parameter}: {val} (Std: {c.standard})")
        checklist_text = "\n".join(checklist_list)

        parts_list = []
        for up in wo.used_parts:
            parts_list.append(f"- {up.part.name} (Qty: {up.quantity_used})")
        parts_text = "\n".join(parts_list)

        data.append({
            'Code': wo.code,
            'Description': wo.description,
            'Status': wo.current_status.name if wo.current_status else '',
            'Priority': wo.priority,
            'Maintenance Type': wo.maintenance_type,
            'Asset Code': wo.asset.code if wo.asset else '',
            'Asset Name': wo.asset.name if wo.asset else '',
            'Site': wo.asset.site.name if wo.asset and wo.asset.site else '',
            'Assigned Technicians': ", ".join([u.name for u in wo.assignees]),
            'Project Code': wo.project_code,
            'Start Date': wo.start_date.strftime('%Y-%m-%d %H:%M') if wo.start_date else '',
            'End Date': wo.end_date.strftime('%Y-%m-%d %H:%M') if wo.end_date else '',
            'Suggested Start': wo.suggested_start_date.strftime('%Y-%m-%d %H:%M') if wo.suggested_start_date else '',
            'Suggested Completion': wo.suggested_completion_date.strftime('%Y-%m-%d %H:%M') if wo.suggested_completion_date else '',
            'Estimated Hours': wo.estimated_hours,
            'Tasklist': wo.tasklist.name if wo.tasklist else '',
            'Checklist': wo.checklist.name if wo.checklist else '',
            'Procedures/Tasks': procedures_text,
            'Checklist Results': checklist_text,
            'Parts Used': parts_text
        })
        
    df = pd.DataFrame(data)
    
    if format == 'excel':
        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            df.to_excel(writer, index=False, sheet_name='Work Orders')
        output.seek(0)
        return send_file(output, as_attachment=True, download_name=f'work_orders_{datetime.now().strftime("%Y%m%d_%H%M")}.xlsx', mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    else:
        output = io.StringIO()
        df.to_csv(output, index=False)
        output.seek(0)
        return Response(output.getvalue(), mimetype="text/csv", headers={"Content-disposition": f"attachment; filename=work_orders_{datetime.now().strftime('%Y%m%d_%H%M')}.csv"})

@work_orders_bp.route('/template')
@login_required
@supervisor_or_admin_required
def download_template():
    columns = [
        'Code', 'Description', 'Asset Name', 'Priority', 'Maintenance Type', 
        'Assignee Name', 'Project Code', 'Tasklist Name', 'Checklist Name',
        'Suggested Start Date', 'Suggested Completion Date'
    ]
    df = pd.DataFrame(columns=columns)
    
    # Add an example row
    example_row = {
        'Code': 'WO-AUTO-001',
        'Description': 'Example maintenance task',
        'Asset Name': 'Chiller 1',
        'Priority': 'Medium',
        'Maintenance Type': 'Corrective',
        'Assignee Name': 'John Doe',
        'Project Code': 'PROJ-01',
        'Tasklist Name': 'General Maintenance',
        'Checklist Name': 'Safety Checklist',
        'Suggested Start Date': datetime.now().strftime('%Y-%m-%d 08:00'),
        'Suggested Completion Date': datetime.now().strftime('%Y-%m-%d 17:00')
    }
    df = pd.concat([df, pd.DataFrame([example_row])], ignore_index=True)
    
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Template')
        
        # --- Add Dropdowns using DataValidation ---
        worksheet = writer.sheets['Template']
        
        # Fetch lookup data (filter by site if Supervisor)
        if current_user.role != 'Admin' and current_user.site_id:
            assets = Asset.query.filter_by(site_id=current_user.site_id).all()
            users = User.query.filter(
                db.or_(User.site_id == current_user.site_id, User.site_id == None)
            ).all()
        else:
            assets = Asset.query.all()
            users = User.query.all()
            
        asset_names = [a.name for a in assets if a.name]
        user_names = [u.name for u in users if u.name]
        
        projects = ProjectCode.query.all()
        project_codes = [p.code for p in projects if p.code]
        
        tasklists = Tasklist.query.all()
        tasklist_names = [t.name for t in tasklists if t.name]
        
        checklists = Checklist.query.all()
        checklist_names = [c.name for c in checklists if c.name]
        
        priorities = ['Highest', 'High', 'Medium', 'Low', 'None']
        maintenance_types = ['Preventive', 'Corrective', 'Inspection', 'Other']
        
        # Create a hidden 'Lookups' sheet
        wb = writer.book
        lookup_sheet = wb.create_sheet('Lookups')
        lookup_sheet.sheet_state = 'hidden'
        
        from openpyxl.utils import get_column_letter
        from openpyxl.worksheet.datavalidation import DataValidation
        
        def write_lookup_col(col_idx, header, data_list):
            lookup_sheet.cell(row=1, column=col_idx, value=header)
            for r_idx, val in enumerate(data_list, start=2):
                lookup_sheet.cell(row=r_idx, column=col_idx, value=val)
            
            col_letter = get_column_letter(col_idx)
            if data_list:
                return f"Lookups!${col_letter}$2:${col_letter}${len(data_list)+1}"
            return None
            
        # Create named range for all assets
        all_assets_range = write_lookup_col(1, 'Asset Names', asset_names)
        if all_assets_range:
            wb.create_named_range('All_Assets', lookup_sheet, all_assets_range.split('!')[1])
            
        priority_range = write_lookup_col(2, 'Priorities', priorities)
        maint_type_range = write_lookup_col(3, 'Maintenance Types', maintenance_types)
        nrp_range = write_lookup_col(4, 'Assignee Names', user_names)
        proj_range = write_lookup_col(5, 'Project Codes', project_codes)
        tasklist_range = write_lookup_col(6, 'Tasklist Names', tasklist_names)
        checklist_range = write_lookup_col(7, 'Checklist Names', checklist_names)
        
        # Create named ranges for each project's assets
        col_offset = 8
        import re
        for project in project_codes:
            proj_assets = [a.name for a in assets if a.project_code == project and a.name]
            if proj_assets:
                rng = write_lookup_col(col_offset, f"{project} Assets", proj_assets)
                # Sanitize name for Excel Named Range (no spaces, hyphens, etc.)
                safe_name = re.sub(r'[^A-Za-z0-9_]', '_', project)
                if not safe_name[0].isalpha() and safe_name[0] != '_':
                    safe_name = '_' + safe_name
                wb.create_named_range(safe_name, lookup_sheet, rng.split('!')[1])
                col_offset += 1
        
        def add_dv(dv_range, cols_str, formula=None):
            f1 = formula if formula else dv_range
            if f1:
                dv = DataValidation(type="list", formula1=f1, allow_blank=True, showErrorMessage=True, errorTitle="Invalid Input", error="Please select a value from the dropdown list.")
                worksheet.add_data_validation(dv)
                dv.add(cols_str)
                
        # Apply to Template sheet columns
        # Asset Name dependent on Project Code (Column G)
        asset_formula = '=IF(ISBLANK(G2), All_Assets, INDIRECT(SUBSTITUTE(SUBSTITUTE(G2, " ", "_"), "-", "_")))'
        add_dv(None, "C2:C1000", formula=asset_formula)
        add_dv(priority_range, "D2:D1000")
        add_dv(maint_type_range, "E2:E1000")
        add_dv(nrp_range, "F2:F1000")
        add_dv(proj_range, "G2:G1000")
        add_dv(tasklist_range, "H2:H1000")
        add_dv(checklist_range, "I2:I1000")
        
        # Add column widths for better UX
        for col in ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H', 'I', 'J', 'K']:
            worksheet.column_dimensions[col].width = 20
        worksheet.column_dimensions['B'].width = 30 # Description
        
    output.seek(0)
    return send_file(output, as_attachment=True, download_name='wo_import_template.xlsx', mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')

@work_orders_bp.route('/import', methods=['POST'])
@login_required
@supervisor_or_admin_required
def import_all():
    if 'file' not in request.files:
        flash('No file part', 'danger')
        return redirect(url_for('work_orders.index'))
        
    file = request.files['file']
    if file.filename == '':
        flash('No selected file', 'danger')
        return redirect(url_for('work_orders.index'))
        
    if file:
        filename = file.filename
        try:
            if filename.endswith('.csv'):
                df = pd.read_csv(file)
            elif filename.endswith('.xlsx') or filename.endswith('.xls'):
                df = pd.read_excel(file)
            else:
                flash('Unsupported file format. Please upload CSV or Excel.', 'danger')
                return redirect(url_for('work_orders.index'))
                
            success_count = 0
            error_count = 0
            errors = []
            
            # Pre-fetch statuses to avoid repeated queries
            draft_status = WorkOrderStatus.query.filter_by(name='Draft').first()
            assigned_status = WorkOrderStatus.query.filter(WorkOrderStatus.name.in_(['Assigned', 'Confirmed'])).first()
            
            for index, row in df.iterrows():
                try:
                    description = row.get('Description')
                    asset_name = row.get('Asset Name')
                    
                    if pd.isna(description) or pd.isna(asset_name):
                        error_count += 1
                        errors.append(f"Row {index+2}: Description or Asset Name is missing.")
                        continue
                    
                    # Find Asset
                    asset = Asset.query.filter_by(name=str(asset_name).strip()).first()
                    if not asset:
                        error_count += 1
                        errors.append(f"Row {index+2}: Asset with name '{asset_name}' not found.")
                        continue
                        
                    # Find Tasklist if name is provided
                    tasklist = None
                    tl_name = row.get('Tasklist Name')
                    if not pd.isna(tl_name):
                        tasklist = Tasklist.query.filter_by(name=str(tl_name).strip()).first()
                        
                    # Find Checklist if name is provided
                    checklist = None
                    cl_name = row.get('Checklist Name')
                    if not pd.isna(cl_name):
                        checklist = Checklist.query.filter_by(name=str(cl_name).strip()).first()
                        
                    # Find Assignee if Name is provided
                    assignee = None
                    assignee_name = row.get('Assignee Name')
                    if not pd.isna(assignee_name):
                        assignee = User.query.filter_by(name=str(assignee_name).strip()).first()
                        
                    # Handle Dates
                    def parse_date(date_val):
                        if pd.isna(date_val): return None
                        if isinstance(date_val, datetime): return date_val
                        try:
                            return datetime.strptime(str(date_val).strip(), '%Y-%m-%d %H:%M')
                        except:
                            try:
                                return datetime.strptime(str(date_val).strip(), '%Y-%m-%d')
                            except:
                                return None
                                
                    suggested_start = parse_date(row.get('Suggested Start Date'))
                    suggested_end = parse_date(row.get('Suggested Completion Date'))
                    
                    # Generate Code if not provided
                    code = row.get('Code')
                    if pd.isna(code) or not str(code).strip():
                        # Use the same logic as create()
                        current_date = datetime.now(timezone.utc).replace(tzinfo=None)
                        month_year = current_date.strftime("%m-%Y")
                        max_id = db.session.query(db.func.max(WorkOrder.id)).scalar() or 0
                        running_number = max_id + 1 + success_count
                        running_number_str = f"{running_number:04d}"
                        project_code = str(row.get('Project Code', 'GEN')) if not pd.isna(row.get('Project Code')) else 'GEN'
                        mt_type = str(row.get('Maintenance Type', 'CM')).strip()
                        type_code = 'CM'
                        if mt_type.lower() in ['preventive', 'pm']:
                            type_code = 'PM'
                        elif mt_type.lower() in ['handling', 'hm']:
                            type_code = 'HM'
                        code = f"WO-{running_number_str}-{project_code}-{type_code}-{month_year}"
                    else:
                        code = str(code).strip()
                        
                    # Check if WO code already exists
                    existing_wo = WorkOrder.query.filter_by(code=code).first()
                    if existing_wo:
                        error_count += 1
                        errors.append(f"Row {index+2}: Work Order with code '{code}' already exists.")
                        continue
                        
                    # Initial Status
                    status_id = assigned_status.id if (assignee and assigned_status) else (draft_status.id if draft_status else None)
                    
                    # Create Work Order
                    new_wo = WorkOrder(
                        code=code,
                        description=str(description),
                        asset_id=asset.id,
                        status_id=status_id,
                        priority=str(row.get('Priority', 'Medium')) if not pd.isna(row.get('Priority')) else 'Medium',
                        maintenance_type=str(row.get('Maintenance Type', 'Corrective')) if not pd.isna(row.get('Maintenance Type')) else 'Corrective',
                        project_code=str(row.get('Project Code')) if not pd.isna(row.get('Project Code')) else None,
                        tasklist_id=tasklist.id if tasklist else None,
                        checklist_id=checklist.id if checklist else None,
                        suggested_start_date=suggested_start,
                        suggested_completion_date=suggested_end
                    )
                    
                    if assignee:
                        new_wo.assignees.append(assignee)
                        
                    db.session.add(new_wo)
                    db.session.flush() # Get the new WO ID
                    
                    # Log Creation
                    l = WorkOrderLog(work_order_id=new_wo.id, user_id=current_user.id, log_text="Work Order created (Imported).")
                    db.session.add(l)
                    db.session.flush() # Generate ID
                    
                    # Copy procedures if tasklist
                    if tasklist:
                        for proc in tasklist.procedures:
                            new_proc = WorkOrderProcedure(
                                work_order_id=new_wo.id,
                                tasklist_name=tasklist.name,
                                name=proc.name,
                                requires_attachment=proc.requires_attachment,
                                min_photos=proc.min_photos,
                                estimated_minutes=proc.estimated_minutes
                            )
                            db.session.add(new_proc)
                            
                    # Copy checklist parameters if checklist
                    if checklist:
                        for param in checklist.parameters:
                            new_chk = WorkOrderChecklistParameter(
                                work_order_id=new_wo.id,
                                checklist_name=checklist.name,
                                parameter=param.parameter,
                                standard=param.standard
                            )
                            db.session.add(new_chk)
                            
                    success_count += 1
                    
                except Exception as e:
                    error_count += 1
                    errors.append(f"Row {index+2}: Internal error - {str(e)}")
                    
            db.session.commit()
            
            if error_count > 0:
                error_msg = f"Import completed with {error_count} errors. {success_count} work orders created."
                if errors:
                    error_msg += " First few errors: " + "; ".join(errors[:3])
                flash(error_msg, 'warning')
            else:
                flash(f'Successfully imported {success_count} work orders.', 'success')
                
        except Exception as e:
            db.session.rollback()
            flash(f'Error processing file: {str(e)}', 'danger')
            
    return redirect(url_for('work_orders.index'))

@work_orders_bp.route('/create', methods=['GET', 'POST'])
@login_required
@supervisor_or_admin_required
def create():
    if request.method == 'POST':
        code = request.form.get('code')
        description = request.form.get('description')
        status_id = request.form.get('status_id')
        priority = request.form.get('priority', 'Medium')
        maintenance_type = request.form.get('maintenance_type', 'Corrective')
        asset_id = request.form.get('asset_id')
        team_id = request.form.get('team_id') or None
        
        assignees = []
        user_ids = set(request.form.getlist('assigned_to'))
        
        # If team is selected, add its members
        if team_id:
            team = db.session.get(Team, team_id)
            if team:
                for u in team.users:
                    user_ids.add(str(u.id))
        
        for uid in user_ids:
            if uid:
                u = db.session.get(User, uid)
                if u:
                    assignees.append(u)
        
        project_code = request.form.get('project_code') or None
            
        tasklist_id = request.form.get('tasklist_id') or None
        if tasklist_id:
            tasklist_id = int(tasklist_id)
            
        checklist_id = request.form.get('checklist_id') or None
        if checklist_id:
            checklist_id = int(checklist_id)
            
        helpdesk_ticket_id = request.form.get('helpdesk_ticket_id') or None
        if helpdesk_ticket_id:
            helpdesk_ticket_id = int(helpdesk_ticket_id)
        
        start_date_str = request.form.get('start_date')
        end_date_str = request.form.get('end_date')
        suggested_start_date_str = request.form.get('suggested_start_date')
        suggested_date_str = request.form.get('suggested_completion_date')

        # Mandatory Field Validation
        project_code = request.form.get('project_code')
        user_site_id = getattr(current_user, 'site_id', None)
        if user_site_id:
            user_site = Site.query.get(user_site_id)
            if user_site and user_site.project_code:
                project_code = user_site.project_code
                
        asset_id = request.form.get('asset_id')
        tasklist_id = request.form.get('tasklist_id')
        checklist_id = request.form.get('checklist_id')
        custom_proc_json = request.form.get('custom_procedures_json', '[]')
        custom_chk_json = request.form.get('custom_checklists_json', '[]')

        with open('debug.txt', 'a') as f:
            f.write(f"POST received. form keys: {list(request.form.keys())}\n")
            f.write(f"proc_json: {custom_proc_json}\n")
            f.write(f"chk_json: {custom_chk_json}\n")
            f.write(f"tasklist_id: {tasklist_id}\n")

        errors = []
        if not asset_id:
            errors.append("Asset selection is mandatory.")
        if not project_code:
            errors.append("Project Code selection is mandatory.")
        
        # Check if tasklist is provided either via template or custom tasks
        has_tasks = False
        if tasklist_id:
            has_tasks = True
        else:
            try:
                if json.loads(custom_proc_json):
                    has_tasks = True
            except: pass
        if not has_tasks:
            errors.append("Tasklist is mandatory (select a template or add custom tasks).")

        # Check if checklist is provided either via template or custom parameters
        has_checklists = False
        if checklist_id:
            has_checklists = True
        else:
            try:
                if json.loads(custom_chk_json):
                    has_checklists = True
            except: pass
        if not has_checklists:
            errors.append("Checklist is mandatory (select a template or add custom parameters).")

        if errors:
            for error in errors:
                flash(error, 'danger')
            return redirect(url_for('work_orders.create', **request.args))
        
        start_date = datetime.strptime(start_date_str, '%Y-%m-%dT%H:%M') if start_date_str else None
        end_date = datetime.strptime(end_date_str, '%Y-%m-%dT%H:%M') if end_date_str else None
        suggested_start_date = datetime.strptime(suggested_start_date_str, '%Y-%m-%dT%H:%M') if suggested_start_date_str else None
        suggested_completion_date = datetime.strptime(suggested_date_str, '%Y-%m-%dT%H:%M') if suggested_date_str else None
        
        estimated_hours = None
        if start_date and end_date:
            estimated_hours = round((end_date - start_date).total_seconds() / 3600, 2)
        
        # basic generation if code is empty
        if not code:
            # Format: WO-[Running Number]-[Project]-[PM/CM]-[Month-Year]
            current_date = datetime.now(timezone.utc).replace(tzinfo=None)
            month_year = current_date.strftime("%m-%Y")
            
            # Running Number
            max_id = db.session.query(db.func.max(WorkOrder.id)).scalar() or 0
            running_number = max_id + 1
            running_number_str = f"{running_number:04d}"
            
            # Project Code
            project_code_gen = project_code if project_code else 'GEN'
            
            # PM/CM/HM Code
            type_code = 'CM'
            if maintenance_type.lower() in ['preventive', 'pm']:
                type_code = 'PM'
            elif maintenance_type.lower() in ['handling', 'hm']:
                type_code = 'HM'
            
            code = f"WO-{running_number_str}-{project_code_gen}-{type_code}-{month_year}"
            
        # State Machine Auto Status handling
        now_dt = datetime.now(timezone.utc).replace(tzinfo=None)
        target_date = suggested_completion_date or suggested_start_date or start_date or now_dt
        is_future_month = (target_date.year > now_dt.year) or (target_date.year == now_dt.year and target_date.month > now_dt.month)

        if is_future_month:
            status = WorkOrderStatus.query.filter_by(name='Draft').first()
        elif assignees:
            status = WorkOrderStatus.query.filter(WorkOrderStatus.name.in_(['Assigned', 'Confirmed'])).first()
        else:
            status = WorkOrderStatus.query.filter_by(name='Draft').first()
            
        if not status:
            target_name = 'Draft' if (is_future_month or not assignees) else 'Assigned'
            status = WorkOrderStatus(name=target_name, control_type='Draft' if target_name == 'Draft' else 'Pending')
            db.session.add(status)
            db.session.flush()
            
        status_id = status.id if status else None
        
        new_wo = WorkOrder(
            code=code, 
            description=description, 
            status_id=status_id, 
            priority=priority, 
            maintenance_type=maintenance_type,
            asset_id=asset_id,
            team_id=team_id,
            assignees=assignees,
            project_code=project_code,
            tasklist_id=tasklist_id,
            checklist_id=checklist_id,
            start_date=start_date,
            end_date=end_date,
            suggested_start_date=suggested_start_date,
            suggested_completion_date=suggested_completion_date,
            estimated_hours=estimated_hours,
            helpdesk_ticket_id=helpdesk_ticket_id
        )
        db.session.add(new_wo)
        db.session.flush()
                    
        import json
        
        # Parse custom procedures
        custom_proc_json = request.form.get('custom_procedures_json', '[]')
        custom_procs_list = []
        try:
            custom_procs_list = json.loads(custom_proc_json)
        except:
            pass
            
        if custom_procs_list:
            for cp in custom_procs_list:
                if isinstance(cp, dict) and cp.get('name'):
                    new_cp = WorkOrderProcedure(
                        work_order=new_wo,
                        tasklist_name='Custom Task',
                        name=cp.get('name'),
                        estimated_minutes=cp.get('estimated_minutes', 0),
                        requires_attachment=cp.get('requires_attachment', False),
                        min_photos=cp.get('min_photos', 0)
                    )
                    db.session.add(new_cp)
        elif tasklist_id:
            # Fallback Native Template Copying
            tasklist = db.session.get(Tasklist, tasklist_id)
            if tasklist:
                for proc in tasklist.procedures:
                    new_proc = WorkOrderProcedure(work_order=new_wo, tasklist_name=tasklist.name, name=proc.name, requires_attachment=proc.requires_attachment, min_photos=proc.min_photos)
                    db.session.add(new_proc)
                    
        # Parse custom checklists
        custom_chk_json = request.form.get('custom_checklists_json', '[]')
        custom_chks_list = []
        try:
            custom_chks_list = json.loads(custom_chk_json)
        except:
            pass
            
        if custom_chks_list:
            for cc in custom_chks_list:
                if isinstance(cc, dict) and cc.get('parameter'):
                    new_cc = WorkOrderChecklistParameter(
                        work_order=new_wo,
                        checklist_name='Custom Parameter',
                        parameter=cc.get('parameter'),
                        standard=cc.get('standard', '')
                    )
                    db.session.add(new_cc)
        elif checklist_id:
            # Fallback Native Template Copying
            checklist = db.session.get(Checklist, checklist_id)
            if checklist:
                for param in checklist.parameters:
                    new_chk = WorkOrderChecklistParameter(work_order=new_wo, checklist_name=checklist.name, parameter=param.parameter, standard=param.standard)
                    db.session.add(new_chk)
                    

                    

        # Process custom parts
        custom_parts_json = request.form.get('custom_parts_json', '[]')
        custom_parts_list = []
        try:
            custom_parts_list = json.loads(custom_parts_json)
        except:
            pass
            
        parts_flash_msg = ""
        if custom_parts_list:
            for cv in custom_parts_list:
                if isinstance(cv, dict):
                    pid = cv.get('part_id')
                    qty = float(cv.get('qty', 0))
                    if pid and qty > 0:
                        part_id = int(pid)
                        # Verify stock — only if asset is loaded and has a site
                        asset_site_id = new_wo.asset.site_id if new_wo.asset else None
                        if asset_site_id:
                            stock = StockLevel.query.filter_by(part_id=part_id, site_id=asset_site_id).first()
                        else:
                            stock = None
                        if stock and stock.qty_on_hand >= qty:
                            stock.qty_on_hand -= qty
                            wo_part = WorkOrderPart.query.filter_by(work_order_id=new_wo.id, part_id=part_id).first()
                            if wo_part:
                                wo_part.quantity_used += qty
                            else:
                                wo_part = WorkOrderPart(work_order_id=new_wo.id, part_id=part_id, quantity_used=qty)
                                db.session.add(wo_part)
                        else:
                            parts_flash_msg = " Some parts were not added due to insufficient stock."
                            
        # Log Creation
        l = WorkOrderLog(work_order_id=new_wo.id, user_id=current_user.id, log_text="Work Order created.")
        db.session.add(l)
                            
        db.session.commit()
        
        # WhatsApp Notification
        for tech in new_wo.assignees:
            if tech.phone_number:
                wa_msg = f"*MAINTENANCE BARU (WO)*\n\nKode: {new_wo.code}\nDeskripsi: {new_wo.description}\nPrioritas: {new_wo.priority}\nAsset: {new_wo.asset.name if new_wo.asset else '-'}\n\nSilakan cek jadwal di sistem."
                send_whatsapp_notification(tech.phone_number, wa_msg)
                
        flash('Work Order created successfully!' + parts_flash_msg, 'success')
        return redirect(url_for('work_orders.index'))
        
    current_date = datetime.now(timezone.utc).replace(tzinfo=None)
    month_year = current_date.strftime("%m-%Y")
    max_id = db.session.query(db.func.max(WorkOrder.id)).scalar() or 0
    running_number = max_id + 1
    running_number_str = f"{running_number:04d}"
        
    user_site_id = getattr(current_user, 'site_id', None)
    if user_site_id:
        assets = Asset.query.filter_by(site_id=user_site_id).all()
        teams = Team.query.filter_by(site_id=user_site_id).all()
        users = User.query.filter(db.or_(User.site_id == user_site_id, User.site_id == None)).all()
        
        user_site = Site.query.get(user_site_id)
        if user_site and user_site.project_code:
            tasklists = Tasklist.query.filter(db.or_(Tasklist.project_code == user_site.project_code, Tasklist.project_code == None, Tasklist.project_code == '')).order_by(Tasklist.name).all()
            checklists = Checklist.query.filter(db.or_(Checklist.project_code == user_site.project_code, Checklist.project_code == None, Checklist.project_code == '')).order_by(Checklist.name).all()
        else:
            tasklists = Tasklist.query.filter(db.or_(Tasklist.project_code == None, Tasklist.project_code == '')).order_by(Tasklist.name).all()
            checklists = Checklist.query.filter(db.or_(Checklist.project_code == None, Checklist.project_code == '')).order_by(Checklist.name).all()
    else:
        assets = Asset.query.all()
        teams = Team.query.all()
        users = User.query.all()
        tasklists = Tasklist.query.order_by(Tasklist.name).all()
        checklists = Checklist.query.order_by(Checklist.name).all()
        
    import json
    assets_json = json.dumps([{'id': a.id, 'name': a.name, 'code': a.code, 'project_code': (a.site.project_code if a.site and a.site.project_code else a.project_code)} for a in assets])

    statuses = WorkOrderStatus.query.all()
    for t in tasklists:
        t.project_code_group = t.project_code if t.project_code else ""
    for c in checklists:
        c.project_code_group = c.project_code if c.project_code else ""
    parts = Part.query.order_by(Part.name).all()
    
    tasklists_json = json.dumps([{'id': t.id, 'name': t.name, 'project_code': t.project_code} for t in tasklists])
    checklists_json = json.dumps([{'id': c.id, 'name': c.name, 'project_code': c.project_code} for c in checklists])
    
    ticket = None
    ticket_id = request.args.get('ticket_id')
    if ticket_id:
        from models import HelpdeskTicket
        ticket = db.session.get(HelpdeskTicket, ticket_id)
    
    return render_template('work_orders/create.html', ticket=ticket, parts=parts, assets=assets, assets_json=assets_json, users=users, teams=teams, statuses=statuses, project_codes=get_all_project_codes(), tasklists=tasklists, tasklists_json=tasklists_json, checklists=checklists, checklists_json=checklists_json, month_year=month_year, running_number_str=running_number_str)

@work_orders_bp.route('/api/tasklist/<int:id>', methods=['GET'])
@login_required
def get_tasklist_api(id):
    tasklist = Tasklist.query.get_or_404(id)
    procs = []
    for p in tasklist.procedures:
        procs.append({
            'name': p.name,
            'estimated_minutes': p.estimated_minutes,
            'requires_attachment': p.requires_attachment
        })
    from flask import jsonify
    return jsonify(procs)

@work_orders_bp.route('/api/checklist/<int:id>', methods=['GET'])
@login_required
def get_checklist_api(id):
    chk = Checklist.query.get_or_404(id)
    params = []
    for p in chk.parameters:
        params.append({
            'parameter': p.parameter,
            'standard': p.standard
        })
    from flask import jsonify
    return jsonify(params)

@work_orders_bp.route('/api/team/<int:id>/members', methods=['GET'])
@login_required
def get_team_members_api(id):
    team = Team.query.get_or_404(id)
    members = []
    for u in team.users:
        members.append({
            'id': u.id,
            'name': u.name,
            'nrp': u.nrp,
            'role': u.role
        })
    from flask import jsonify
    return jsonify(members)

@work_orders_bp.route('/<int:id>/edit', methods=['GET', 'POST'])
@login_required
def edit(id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        flash('Access denied. You can only access work orders assigned to you.', 'danger')
        return redirect(url_for('work_orders.index'))
        
    # Self-healing logic for WOs created without a status
    if not wo.status_id:
        target_name = 'Assigned' if wo.assignees else 'Draft'
        status = WorkOrderStatus.query.filter(WorkOrderStatus.name.in_(['Assigned', 'Confirmed'])).first() if target_name == 'Assigned' else WorkOrderStatus.query.filter_by(name='Draft').first()
        if not status:
            status = WorkOrderStatus(name=target_name, control_type='Draft' if target_name == 'Draft' else 'Pending')
            db.session.add(status)
            db.session.flush()
        wo.status_id = status.id
        db.session.commit()
        
    if request.method == 'POST':
        # Prevent edits if WO is in terminal state, but allow signature updates
        if wo.current_status and wo.current_status.name in ['Solved', 'Completed', 'Closed', 'Incompleted']:
            if not request.form.get('customer_signature') and not request.form.get('technician_signature') and not request.form.get('customer_name'):
                flash('This Work Order is already solved/completed and cannot be modified.', 'danger')
                return redirect(url_for('work_orders.edit', id=wo.id))
            else:
                # Process only signature fields and return early
                signature_added = False
                if request.form.get('customer_name') is not None:
                    wo.customer_name = request.form.get('customer_name')
                if request.form.get('customer_title') is not None:
                    wo.customer_title = request.form.get('customer_title')
                    
                c_sig = request.form.get('customer_signature')
                if c_sig and c_sig.startswith('data:image'):
                    wo.customer_signature = c_sig
                    signature_added = True
                    
                t_sig = request.form.get('technician_signature')
                if t_sig and t_sig.startswith('data:image'):
                    wo.technician_signature = t_sig
                    signature_added = True
                    
                if signature_added:
                    from models import DigitalSignature
                    existing_sig = DigitalSignature.query.filter_by(work_order_id=wo.id).first()
                    if not existing_sig:
                        signer = (wo.customer_name or current_user.name)
                        title = (wo.customer_title or 'Technician')
                        import uuid, hashlib, time
                        sig_id = str(uuid.uuid4())
                        ts = str(time.time())
                        data_to_hash = f"{wo.code}|{signer}|{ts}"
                        doc_hash = hashlib.sha256(data_to_hash.encode('utf-8')).hexdigest()
                        new_sig = DigitalSignature(
                            id=sig_id, work_order_id=wo.id, signer_name=signer, 
                            signer_title=title, signed_at=datetime.utcnow(), 
                            document_hash=doc_hash, status='Valid'
                        )
                        db.session.add(new_sig)
                        
                db.session.commit()
                flash('Signatures saved successfully.', 'success')
                return redirect(url_for('work_orders.edit', id=wo.id))

        desc = request.form.get('description')
        if desc is not None:
            wo.description = desc
            
        prio = request.form.get('priority')
        if prio is not None:
            wo.priority = prio
            
        m_type = request.form.get('maintenance_type')
        if m_type is not None:
            wo.maintenance_type = m_type
            
        team_id = request.form.get('team_id')
        if team_id is not None:
            wo.team_id = team_id or None
        
        new_asset = request.form.get('asset_id')
        if new_asset and (not wo.current_status or wo.current_status.name == 'Draft'):
            wo.asset_id = new_asset
            
        new_tl_id = request.form.get('tasklist_id') or None
        if new_tl_id:
            new_tl_id = int(new_tl_id)
            
        new_cl_id = request.form.get('checklist_id') or None
        if new_cl_id:
            new_cl_id = int(new_cl_id)
            
        # Handle Tasklist Change
        if wo.tasklist_id != new_tl_id:
            wo.tasklist_id = new_tl_id
            for proc in wo.procedures:
                db.session.delete(proc)
            if new_tl_id:
                curr_tl = db.session.get(Tasklist, new_tl_id)
                if curr_tl:
                    for proc in curr_tl.procedures:
                        new_proc = WorkOrderProcedure(work_order=wo, tasklist_name=curr_tl.name, name=proc.name, requires_attachment=proc.requires_attachment, min_photos=proc.min_photos, estimated_minutes=proc.estimated_minutes)
                        db.session.add(new_proc)
                        
        # Handle Checklist Change
        if wo.checklist_id != new_cl_id:
            wo.checklist_id = new_cl_id
            for chk in wo.checklist_parameters:
                db.session.delete(chk)
            if new_cl_id:
                curr_cl = db.session.get(Checklist, new_cl_id)
                if curr_cl:
                    for param in curr_cl.parameters:
                        new_chk = WorkOrderChecklistParameter(work_order=wo, checklist_name=curr_cl.name, parameter=param.parameter, standard=param.standard)
                        db.session.add(new_chk)
        
        assigned = request.form.getlist('assigned_to')
        if 'assigned_to' in request.form:
            new_assignees = []
            for uid in assigned:
                if uid:
                    u = db.session.get(User, uid)
                    if u:
                        new_assignees.append(u)
            
            # Auto Upgrade Draft to Assigned if assigning a technician
            if new_assignees and wo.current_status and wo.current_status.name == 'Draft':
                assigned = WorkOrderStatus.query.filter(WorkOrderStatus.name.in_(['Assigned', 'Confirmed'])).first()
                if assigned:
                    wo.status_id = assigned.id
                    l = WorkOrderLog(work_order_id=wo.id, user_id=current_user.id, log_text="Status changed to Assigned (Auto-assigned).")
                    db.session.add(l)
                    
            # Auto Downgrade to Draft if unassigning while Assigned
            if not new_assignees and wo.current_status and wo.current_status.name in ['Assigned', 'Confirmed']:
                draft_status = WorkOrderStatus.query.filter_by(name='Draft').first()
                if draft_status:
                    wo.status_id = draft_status.id
                    l = WorkOrderLog(work_order_id=wo.id, user_id=current_user.id, log_text="Status changed to Draft (Unassigned).")
                    db.session.add(l)
                    
            wo.assignees = new_assignees
            
        p_code = request.form.get('project_code')
        user_site_id = getattr(current_user, 'site_id', None)
        if user_site_id:
            user_site = Site.query.get(user_site_id)
            if user_site and user_site.project_code:
                p_code = user_site.project_code
                
        if p_code is not None:
            wo.project_code = p_code or None
        
        if 'start_date' in request.form:
            start_date_str = request.form.get('start_date')
            wo.start_date = datetime.strptime(start_date_str, '%Y-%m-%dT%H:%M') if start_date_str else None
            
        if 'end_date' in request.form:
            end_date_str = request.form.get('end_date')
            wo.end_date = datetime.strptime(end_date_str, '%Y-%m-%dT%H:%M') if end_date_str else None
            
        if 'suggested_start_date' in request.form:
            suggested_start_date_str = request.form.get('suggested_start_date')
            wo.suggested_start_date = datetime.strptime(suggested_start_date_str, '%Y-%m-%dT%H:%M') if suggested_start_date_str else None
            
        if 'suggested_completion_date' in request.form:
            suggested_date_str = request.form.get('suggested_completion_date')
            wo.suggested_completion_date = datetime.strptime(suggested_date_str, '%Y-%m-%dT%H:%M') if suggested_date_str else None
        
        if wo.start_date and wo.end_date:
            wo.estimated_hours = round((wo.end_date - wo.start_date).total_seconds() / 3600, 2)
        else:
            wo.estimated_hours = None
            
        # Signatures
        signature_added = False
        if request.form.get('customer_name') is not None:
            wo.customer_name = request.form.get('customer_name')
        if request.form.get('customer_title') is not None:
            wo.customer_title = request.form.get('customer_title')
            
        c_sig = request.form.get('customer_signature')
        if c_sig and c_sig.startswith('data:image'):
            wo.customer_signature = c_sig
            signature_added = True
            
        t_sig = request.form.get('technician_signature')
        if t_sig and t_sig.startswith('data:image'):
            wo.technician_signature = t_sig
            signature_added = True
            
        if signature_added:
            from models import DigitalSignature
            existing_sig = DigitalSignature.query.filter_by(work_order_id=wo.id).first()
            if not existing_sig:
                signer = (wo.customer_name or current_user.name)
                title = (wo.customer_title or 'Technician')
                import uuid, hashlib, time
                sig_id = str(uuid.uuid4())
                ts = str(time.time())
                data_to_hash = f"{wo.code}|{signer}|{ts}"
                doc_hash = hashlib.sha256(data_to_hash.encode('utf-8')).hexdigest()
                new_sig = DigitalSignature(
                    id=sig_id,
                    work_order_id=wo.id,
                    signer_name=signer,
                    signer_title=title,
                    signed_at=datetime.utcnow(),
                    document_hash=doc_hash,
                    status='Valid'
                )
                db.session.add(new_sig)
            
        # File Upload Handling
        if 'wo_attachment' in request.files:
            file = request.files['wo_attachment']
            if file and file.filename != '':
                filename = secure_filename(file.filename)
                
                # Ensure uploads dir exists
                uploads_dir = os.path.join(current_app.root_path, 'static', 'uploads')
                os.makedirs(uploads_dir, exist_ok=True)
                
                # Save file
                file_path_sys = os.path.join(uploads_dir, filename)
                file.save(file_path_sys)
                
                # DB relative path for URL generation
                db_path = f'uploads/{filename}'
                
                # Create attachment record
                from models import WorkOrderAttachment
                new_att = WorkOrderAttachment(work_order_id=wo.id, file_name=filename, file_path=db_path)
                db.session.add(new_att)
        
        db.session.commit()
        flash('Work Order updated successfully!', 'success')
        return redirect(url_for('work_orders.edit', id=wo.id))
        
    if current_user.site_id:
        assets = Asset.query.filter_by(site_id=current_user.site_id).all()
    else:
        assets = Asset.query.all()
        
    import json
    assets_json = json.dumps([{'id': a.id, 'name': a.name, 'code': a.code, 'project_code': (a.site.project_code if a.site and a.site.project_code else a.project_code)} for a in assets])

    user_site_id = getattr(current_user, 'site_id', None)
    if user_site_id:
        users = User.query.filter(db.or_(User.site_id == user_site_id, User.site_id == None)).all()
        user_site = Site.query.get(user_site_id)
        if user_site and user_site.project_code:
            tasklists = Tasklist.query.filter(db.or_(Tasklist.project_code == user_site.project_code, Tasklist.project_code == None, Tasklist.project_code == '')).order_by(Tasklist.name).all()
            checklists = Checklist.query.filter(db.or_(Checklist.project_code == user_site.project_code, Checklist.project_code == None, Checklist.project_code == '')).order_by(Checklist.name).all()
        else:
            tasklists = Tasklist.query.filter(db.or_(Tasklist.project_code == None, Tasklist.project_code == '')).order_by(Tasklist.name).all()
            checklists = Checklist.query.filter(db.or_(Checklist.project_code == None, Checklist.project_code == '')).order_by(Checklist.name).all()
    else:
        users = User.query.all()
        tasklists = Tasklist.query.order_by(Tasklist.name).all()
        checklists = Checklist.query.order_by(Checklist.name).all()
        
    statuses = WorkOrderStatus.query.all()
    for t in tasklists:
        t.project_code_group = t.project_code if t.project_code else ""
    for c in checklists:
        c.project_code_group = c.project_code if c.project_code else ""
    
    tasklists_json = json.dumps([{'id': t.id, 'name': t.name, 'project_code': t.project_code} for t in tasklists])
    checklists_json = json.dumps([{'id': c.id, 'name': c.name, 'project_code': c.project_code} for c in checklists])
    users_json = json.dumps([{'id': u.id, 'name': u.name, 'role': u.role, 'project_code': u.site.project_code if u.site and u.site.project_code else None} for u in users])
    
    if wo.asset and wo.asset.site_id:
        available_stock = StockLevel.query.filter_by(site_id=wo.asset.site_id).all()
        teams = Team.query.filter_by(site_id=wo.asset.site_id).all()
    else:
        available_stock = []
        teams = Team.query.all()
        
    return render_template('work_orders/edit.html', wo=wo, assets=assets, assets_json=assets_json, users=users, users_json=users_json, teams=teams, available_stock=available_stock, statuses=statuses, project_codes=get_all_project_codes(), tasklists=tasklists, tasklists_json=tasklists_json, checklists=checklists, checklists_json=checklists_json)


@work_orders_bp.route('/<int:id>/export_pdf', methods=['GET'])
@login_required
def export_pdf(id):
    import os
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        flash('Access denied.', 'danger')
        return redirect(url_for('work_orders.index'))
        
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    
    # --- 0. COVER PAGE ---
    pdf.add_page()
    
    # Background Image
    title_image_path = os.path.join(current_app.root_path, 'static', 'images', 'judul.png')
    if os.path.exists(title_image_path):
        pdf.image(title_image_path, x=0, y=0, w=pdf.w, h=pdf.h)
    
    # Get Site Name and Period
    site_name = wo.asset.site.name if wo.asset and wo.asset.site else "Unknown Site"
    wo_date = wo.suggested_start_date or wo.start_date or datetime.now()
    month_map = {
        1: "JANUARI", 2: "FEBRUARI", 3: "MARET", 4: "APRIL", 5: "MEI", 6: "JUNI",
        7: "JULI", 8: "AGUSTUS", 9: "SEPTEMBER", 10: "OKTOBER", 11: "NOVEMBER", 12: "DESEMBER"
    }
    periode_str = f"{month_map.get(wo_date.month, '')} {wo_date.year}"
    
    mt_type = wo.maintenance_type.upper() if wo.maintenance_type else ""
    if "MAINTENANCE" not in mt_type:
        mt_type += " MAINTENANCE"
    
    line1 = f"LAPORAN WORK ORDER"
    line1_sub = mt_type
    line2 = f"{site_name.upper()}"
    line3 = f"{wo.asset.name.upper() if wo.asset else ''}"
    line4 = f"{wo.code}"
    
    # Positioning text at bottom left
    pdf.set_y(-105) # Moved slightly up to accommodate extra line
    pdf.set_left_margin(20)
    
    pdf.set_font('helvetica', 'B', 24) # Reduced size slightly for line1
    pdf.cell(0, 12, line1, 0, 1, 'L')
    pdf.set_font('helvetica', 'B', 24)
    pdf.cell(0, 12, line1_sub, 0, 1, 'L')
    
    pdf.set_font('helvetica', 'B', 18) # Reduced size for others
    pdf.cell(0, 10, line2, 0, 1, 'L')
    
    pdf.set_font('helvetica', 'B', 18)
    pdf.cell(0, 10, line3, 0, 1, 'L')
    
    pdf.set_font('helvetica', 'B', 18)
    pdf.cell(0, 10, line4, 0, 1, 'L')
    
    # Reset for main report
    pdf.set_left_margin(10)
    pdf.add_page()
    # --- END COVER PAGE ---
    
    try:
        from datetime import datetime, timezone
        pdf.set_font("helvetica", style="B", size=10)
    except:
        pdf.set_font("Arial", style="B", size=10)
        
    start_x = pdf.get_x()
    start_y = pdf.get_y()
    
    # 1. NEW HEADER
    pdf.cell(35, 24, "", border=1) # Logo area box
    
    logo_path = os.path.join(current_app.root_path, 'static', 'images', 'Logo Jaya Teknik.png')
    if os.path.exists(logo_path):
        try:
            pdf.image(logo_path, x=start_x + 2, y=start_y + 2, w=31)
        except:
            pass
    
    # Title
    pdf.set_font("helvetica", style="B", size=14)
    pdf.set_xy(start_x + 35, start_y)
    pdf.cell(105, 24, "", border=1)
   
    pdf.set_xy(start_x + 35, start_y + 12)
    pdf.cell(105, 8, "REPORT WORK ORDER", align="C")
    
    # Info box right
    pdf.set_font("helvetica", size=8)
    pdf.set_xy(start_x + 140, start_y)
    pdf.cell(20, 6, "No.Dok", border=1)
    pdf.cell(30, 6, "RP - JT - 26", border=1)
    
    pdf.set_xy(start_x + 140, start_y + 6)
    pdf.cell(20, 6, "Ref.", border=1)
    pdf.cell(30, 6, "7; 7.1.3 & 7.1.4", border=1)
    
    pdf.set_xy(start_x + 140, start_y + 12)
    pdf.cell(20, 6, "Rev.", border=1)
    pdf.cell(30, 6, "Original", border=1)
    
    pdf.set_xy(start_x + 140, start_y + 18)
    pdf.cell(20, 6, "Tanggal", border=1)
    printed_date = datetime.now().strftime("%d %b %Y")
    pdf.cell(30, 6, printed_date, border=1)
    
    # Details Body
    pdf.set_xy(start_x, start_y + 24)
    pdf.set_font("helvetica", size=9)
    pdf.cell(190, 24, "", border=1)
    pdf.set_xy(start_x, start_y + 24)
    
    def fmt_dt(dt, convert=True):
        from datetime import timedelta
        if dt:
            if convert:
                dt_wib = dt + timedelta(hours=7)
                return dt_wib.strftime('%d-%m-%Y %H:%M')
            return dt.strftime('%d-%m-%Y %H:%M')
        return '-'
        
    form_val = getattr(wo, 'maintenance_type', 'SERVICE')
    nomor_val = wo.code or '-'
    asset_val = f"{wo.asset.name} ({wo.asset.code})" if wo.asset else '-'
    resp_val = ", ".join([a.name for a in wo.assignees]) if getattr(wo, 'assignees', None) else "N/A"
    
    proj_val = (wo.asset.project_code if wo.asset and wo.asset.project_code else wo.project_code) or '-'
    
    loc_val = '-'
    if wo.asset and hasattr(wo.asset, 'location') and wo.asset.location:
        loc_val = wo.asset.location.name if hasattr(wo.asset.location, 'name') else str(wo.asset.location)
    
    sched_val = f"{fmt_dt(wo.suggested_start_date)} s/d {fmt_dt(wo.suggested_completion_date)}"
    act_val = f"{fmt_dt(wo.start_date)} s/d {fmt_dt(wo.end_date)}"
    
    # Line 1
    pdf.cell(20, 6, "Form", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(form_val)[:35], border=0)
    pdf.cell(20, 6, "Project", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(proj_val)[:35], border=0, ln=1)
    
    # Line 2
    pdf.cell(20, 6, "Nomor", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(nomor_val)[:35], border=0)
    pdf.cell(20, 6, "Location", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(loc_val)[:35], border=0, ln=1)
    
    # Line 3
    pdf.cell(20, 6, "Asset", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(asset_val)[:35], border=0)
    pdf.cell(20, 6, "Schedule", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(sched_val)[:40], border=0, ln=1)
    
    # Line 4
    pdf.cell(20, 6, "Responsible", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(resp_val)[:35], border=0)
    pdf.cell(20, 6, "Actual Date", border=0)
    pdf.cell(5, 6, ":", border=0)
    pdf.cell(70, 6, str(act_val)[:40], border=0, ln=1)
    
    pdf.ln(5)
    
    # 1.5 TECHNICIAN (Restored)
    pdf.set_fill_color(220, 220, 220)
    pdf.set_font("helvetica", style="B", size=10)
    pdf.cell(190, 6, "TECHNICIAN", border=1, align="C", fill=True, ln=1)
    pdf.cell(190, 6, "Name", border=1, align="C", fill=True, ln=1)
    pdf.set_font("helvetica", size=9)
    pdf.cell(190, 6, str(resp_val), border=1, align="L", ln=1)
    pdf.ln(5)
    
    # 2. CHECKING REPORT
    pdf.set_fill_color(220, 220, 220)
    pdf.set_font("helvetica", style="B", size=10)
    pdf.cell(190, 6, "CHECKING REPORT", border=1, align="C", fill=True, ln=1)
    
    pdf.set_font("helvetica", style="B", size=8)
    pdf.cell(75, 8, "Description (Unit Check)", border=1, align="C", fill=True)
    pdf.cell(25, 8, "Standard", border=1, align="C", fill=True)
    pdf.cell(25, 8, "Actual", border=1, align="C", fill=True)
    pdf.cell(25, 8, "Check", border=1, align="C", fill=True)
    pdf.cell(40, 8, "Note", border=1, align="C", fill=True, ln=1)
    pdf.set_font("helvetica", size=8)
    for param in wo.checklist_parameters:
        desc = (param.parameter[:40] + '...') if len(param.parameter) > 42 else param.parameter
        pdf.cell(75, 6, str(desc), border=1)
        pdf.cell(25, 6, str(param.standard)[:15], border=1, align="C")
        
        actual_val = str(param.value) if param.value else "-"
        pdf.cell(25, 6, actual_val[:15], border=1, align="C")
        
        check_val = "v" if actual_val != "-" else "-"
        pdf.cell(25, 6, "OK" if check_val == "v" else "-", border=1, align="C")
        
        note_val = str(param.note) if param.note else ""
        pdf.cell(40, 6, note_val[:20], border=1, ln=1)
    
    pdf.ln(5)
    
    # 3. TASKLIST REPORT
    pdf.set_font("helvetica", style="B", size=10)
    pdf.cell(190, 6, "TASKLIST REPORT", border=1, align="C", fill=True, ln=1)
    
    pdf.set_font("helvetica", style="B", size=8)
    pdf.cell(70, 8, "Step Name", border=1, align="C", fill=True)
    pdf.cell(90, 8, "Description", border=1, align="C", fill=True)
    pdf.cell(30, 8, "State", border=1, align="C", fill=True, ln=1)
    pdf.set_font("helvetica", size=8)
    for proc in wo.procedures:
        step_name = proc.name or ''
        desc = getattr(proc, 'description', '') or ''
        
        # Estimate number of lines based on string width
        w1, w2 = 68, 88 # Effective widths for text wrapping
        
        lines_step = max(1, int(pdf.get_string_width(step_name) / w1) + 1) + step_name.count('\n') if step_name else 1
        lines_desc = max(1, int(pdf.get_string_width(desc) / w2) + 1) + desc.count('\n') if desc else 1
        
        max_lines = max(lines_step, lines_desc)
        row_height = max_lines * 6
        
        # Check page break
        if pdf.get_y() + row_height > 270:
            pdf.add_page()
            
        x_start = pdf.get_x()
        y_start = pdf.get_y()
        
        # Draw borders for the row
        pdf.rect(x_start, y_start, 70, row_height)
        pdf.rect(x_start + 70, y_start, 90, row_height)
        pdf.rect(x_start + 160, y_start, 30, row_height)
        
        # Print Step Name
        pdf.set_xy(x_start, y_start)
        pdf.multi_cell(70, 6, step_name, border=0)
        
        # Print Description
        pdf.set_xy(x_start + 70, y_start)
        pdf.multi_cell(90, 6, desc, border=0)
        
        # Print State
        pdf.set_xy(x_start + 160, y_start)
        pdf.cell(30, row_height, "Done" if getattr(proc, 'is_completed', False) else "Pending", border=0, align="C")
        
        pdf.set_y(y_start + row_height)
        
    pdf.ln(5)
    

    
    # 3.6 SIGNATURES
    pdf.ln(5)
    
    first_tech = wo.assignees[0] if getattr(wo, 'assignees', []) else None
    tech_name = first_tech.name if first_tech else "....."
    tech_role = first_tech.role if first_tech else "....."
    site_name = wo.asset.site.name.upper() if wo.asset and getattr(wo.asset, 'site', None) else "....."
    
    if pdf.get_y() > 240:
        pdf.add_page()
        
    sig_y = pdf.get_y()
    
    # Helper to decode signature
    import tempfile
    import base64
    def save_temp_image(b64_str):
        if not b64_str or ',' not in b64_str: return None
        try:
            _, encoded = b64_str.split(",", 1)
            data = base64.b64decode(encoded)
            fd, temp_path = tempfile.mkstemp(suffix='.png')
            with os.fdopen(fd, 'wb') as f:
                f.write(data)
            return temp_path
        except:
            return None

    c_name_print = getattr(wo, 'customer_name', None) or "....."
    c_title_print = getattr(wo, 'customer_title', None) or "....."

    # Fetch Digital Signature first
    from models import DigitalSignature
    ds = DigitalSignature.query.filter_by(work_order_id=wo.id).first()
    
    # Left Column (Pihak Pertama)
    pdf.set_xy(15, sig_y)
    pdf.set_font("helvetica", style="B", size=9)
    pdf.cell(90, 5, "PIHAK PERTAMA,", ln=1)
    pdf.set_x(15)
    pdf.cell(90, 5, site_name, ln=1)
    
    if not ds:
        c_img_path = save_temp_image(getattr(wo, 'customer_signature', None))
        if c_img_path:
            try:
                pdf.image(c_img_path, x=15, y=sig_y + 12, h=16)
                os.remove(c_img_path)
            except: pass

    pdf.set_xy(15, sig_y + 35)
    pdf.set_font("helvetica", size=9)
    pdf.cell(90, 5, f"Nama : {c_name_print}", ln=1)
    pdf.set_x(15)
    pdf.cell(90, 5, f"Jabatan : {c_title_print}", ln=1)
    
    # Right Column (Pihak Kedua)
    pdf.set_xy(110, sig_y)
    pdf.set_font("helvetica", style="B", size=9)
    pdf.cell(85, 5, "PIHAK KEDUA,", ln=1)
    pdf.set_x(110)
    pdf.cell(85, 5, "PT JAYA TEKNIK INDONESIA", ln=1)
    
    if not ds:
        t_img_path = save_temp_image(getattr(wo, 'technician_signature', None))
        if t_img_path:
            try:
                pdf.image(t_img_path, x=110, y=sig_y + 12, h=16)
                os.remove(t_img_path)
            except: pass

    pdf.set_xy(110, sig_y + 35)
    pdf.set_font("helvetica", size=9)
    pdf.cell(85, 5, f"Nama : {tech_name}", ln=1)
    pdf.set_x(110)
    pdf.cell(85, 5, f"Jabatan : {tech_role}", ln=1)
    
    # Center Column (Digital Signature QR)
    if ds:
        import qrcode
        import tempfile
        from flask import request
        
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=1,
        )
        verify_url = f"{request.url_root.rstrip('/')}/verify/doc?id={ds.id}"
        qr.add_data(verify_url)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        
        fd, temp_path = tempfile.mkstemp(suffix='.png')
        with os.fdopen(fd, 'wb') as f:
            img.save(f)
            
        pdf.image(temp_path, x=135, y=sig_y + 10, h=18)
        os.remove(temp_path)
        
        pdf.set_xy(125, sig_y + 29)
        pdf.set_font("helvetica", size=6, style="B")
        pdf.set_text_color(0, 150, 0) # Green text for verification
        pdf.cell(38, 3, "DIGITALLY SIGNED", ln=1, align="C")
        pdf.set_xy(125, sig_y + 32)
        pdf.cell(38, 3, "& VERIFIED", ln=1, align="C")
        pdf.set_text_color(0, 0, 0)
    
    pdf.set_y(sig_y + 45)
       # 4. PHOTOS
    images = []
    for proc in wo.procedures:
        if proc.attachment_path:
             for path in proc.attachment_path.split('|'):
                 if path:
                     images.append({"caption": f"Lampiran {proc.name[:50]}", "path": path})
    for att in wo.attachments:
         images.append({"caption": f"Lampiran {att.file_name}", "path": att.file_path})
         
    box_width = 190 / 3
    box_height = 95
    img_width = 57
    
    if images:
         pdf.add_page()
         pdf.set_fill_color(220, 220, 220)
         pdf.set_font("helvetica", style="B", size=10)
         pdf.cell(190, 6, "LAMPIRAN FOTO", border=1, align="C", fill=True, ln=1)
         
         y_cursor = pdf.get_y()
         for i in range(0, len(images), 3):
             row_images = images[i:i+3]
             n_imgs = len(row_images)
             
             if y_cursor + box_height > 280:
                 pdf.add_page()
                 y_cursor = pdf.get_y()
                 
             # Draw full-width outer border for the row
             pdf.set_xy(10, y_cursor)
             pdf.cell(190, box_height, "", border=1)
             
             # Group contiguous images with the same caption in this row
             groups = []
             for img in row_images:
                 if not groups or groups[-1]['cap'] != img['caption']:
                     groups.append({'cap': img['caption'], 'count': 1})
                 else:
                     groups[-1]['count'] += 1
                     
             # Draw captions
             pdf.set_font("helvetica", style="B", size=7)
             cell_w = 190 / n_imgs
             current_x = 10
             for g in groups:
                 group_w = g['count'] * cell_w
                 cap = g['cap']
                 # Calculate approx max chars based on width
                 max_chars = int((group_w / 190) * 250)
                 display_cap = (cap[:max_chars] + "...") if len(cap) > max_chars else cap
                 
                 pdf.set_xy(current_x, y_cursor + 1)
                 pdf.multi_cell(group_w, 3.5, display_cap, align="C")
                 current_x += group_w
                 
             # Draw images
             for j, img_dict in enumerate(row_images):
                 x_center = 10 + j * cell_w + cell_w / 2
                 x_pos = x_center - (img_width / 2)
                 
                 full_path = os.path.join(current_app.root_path, 'static', img_dict['path'])
                 if os.path.exists(full_path):
                     try:
                         from PIL import ImageOps
                         with Image.open(full_path) as im:
                             im = ImageOps.exif_transpose(im)
                             w, h = im.size
                             
                             if w > h:
                                 im = im.rotate(270, expand=True)
                                 w, h = im.size
                                 
                             target_aspect = 3/4
                             aspect = w/h
                             if aspect > target_aspect:
                                 new_w = int(h * target_aspect)
                                 left = (w - new_w) / 2
                                 im = im.crop((left, 0, left + new_w, h))
                             else:
                                 new_h = int(w / target_aspect)
                                 top = (h - new_h) / 2
                                 im = im.crop((0, top, w, top + new_h))
                             
                             im = im.resize((300, 400))
                             
                             tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
                             if im.mode != 'RGB':
                                 im = im.convert('RGB')
                             im.save(tmp.name, 'JPEG')
                             tmp.close()
                             
                             pdf.image(tmp.name, x=x_pos, y=y_cursor + 13, w=img_width)
                             os.unlink(tmp.name)
                     except Exception as e:
                         print(f"Error processing image {full_path}: {e}")
             
             y_cursor += box_height
    
    try:
        out = pdf.output(dest='S')
    except TypeError:
        out = pdf.output()
    pdf_out = out.encode('latin-1') if isinstance(out, str) else bytes(out)
    
    if ds and not ds.pdf_path:
        sig_dir = os.path.join(current_app.root_path, 'static', 'signatures')
        os.makedirs(sig_dir, exist_ok=True)
        pdf_filename = f"signed_{wo.code}_{ds.id[:8]}.pdf".replace('/', '_')
        pdf_filepath = os.path.join(sig_dir, pdf_filename)
        with open(pdf_filepath, 'wb') as f:
            f.write(pdf_out)
        ds.pdf_path = f"signatures/{pdf_filename}"
        db.session.commit()

    return Response(pdf_out, mimetype='application/pdf', headers={'Content-Disposition': f'attachment; filename="Report_WO_{wo.code}.pdf"'})

@work_orders_bp.route('/<int:id>/status', methods=['POST'])
@login_required
def change_status(id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        flash('Access denied.', 'danger')
        return redirect(url_for('work_orders.index'))
        
    action = request.form.get('action') # 'start', 'hold', 'resume', 'solve', 'incomplete', 'complete'
    reason = request.form.get('reason') # for hold or incomplete
    
    # Check permissions logic
    from flask import abort
    if current_user.role == 'Technician' and action in ['incomplete', 'complete']:
        abort(403)
        
    action_map = {
        'start': 'In Progress',
        'hold': 'Hold',
        'resume': 'In Progress',
        'solve': 'Solved',
        'incomplete': 'Incompleted',
        'complete': 'Completed',
        'cancel_to_draft': 'Draft'
    }
    
    target_status_name = action_map.get(action)
    if not target_status_name:
        flash('Invalid status action.', 'danger')
        return redirect(url_for('work_orders.edit', id=wo.id))
        
    target_status = WorkOrderStatus.query.filter_by(name=target_status_name).first()
    if target_status:
        wo.status_id = target_status.id
        
        # Start Progress Logic -> Populate start_date, Clear end_date
        if action in ['start', 'resume']:
            if action == 'start' and not wo.start_date:
                wo.start_date = datetime.now(timezone.utc).replace(tzinfo=None)
            wo.end_date = None
            wo.estimated_hours = None
            
        # Complete / Solve Logic -> Populate end_date and Verify
        if action in ['solve', 'complete']:
            # Validation: Tasklist Procedures and Checklists
            incomplete_procs = []
            for proc in wo.procedures:
                num_photos = 0
                if proc.attachment_path:
                    num_photos = len([p for p in proc.attachment_path.split('|') if p])
                
                if not proc.is_completed:
                    incomplete_procs.append(proc.name)
                elif proc.requires_attachment and not proc.attachment_path:
                    incomplete_procs.append(f"{proc.name} (Missing Photo)")
                elif proc.min_photos > 0 and num_photos < proc.min_photos:
                    incomplete_procs.append(f"{proc.name} (Photos: {num_photos}/{proc.min_photos})")
                    
            if incomplete_procs:    
                flash(f'Cannot {action}: Incomplete procedures or missing photos: {", ".join(incomplete_procs)}', 'danger')
                return redirect(url_for('work_orders.edit', id=wo.id))
                
            missing_chks = [c.parameter for c in wo.checklist_parameters if not c.value]
            if missing_chks:
                flash(f'Cannot {action}: Missing values for checklist parameters: {", ".join(missing_chks)}', 'danger')
                return redirect(url_for('work_orders.edit', id=wo.id))
            
            wo.end_date = datetime.now(timezone.utc).replace(tzinfo=None)
                
            if wo.start_date and wo.end_date:
                wo.estimated_hours = round((wo.end_date - wo.start_date).total_seconds() / 3600, 2)
                
            # Digital Signature Logic
            from models import DigitalSignature
            if getattr(wo, 'customer_signature', None) or getattr(wo, 'technician_signature', None):
                # Check if already has signature
                existing_sig = DigitalSignature.query.filter_by(work_order_id=wo.id).first()
                if not existing_sig:
                    import uuid
                    import hashlib
                    import time
                    
                    signer = (wo.customer_name or current_user.name)
                    title = (wo.customer_title or 'Technician')
                    
                    sig_id = str(uuid.uuid4())
                    ts = str(time.time())
                    
                    data_to_hash = f"{wo.code}|{signer}|{ts}"
                    doc_hash = hashlib.sha256(data_to_hash.encode('utf-8')).hexdigest()
                    
                    new_sig = DigitalSignature(
                        id=sig_id,
                        work_order_id=wo.id,
                        signer_name=signer,
                        signer_title=title,
                        signed_at=datetime.utcnow(),
                        document_hash=doc_hash,
                        status='Valid'
                    )
                    db.session.add(new_sig)
            
        # Log Status Change
        log_title = f"Status changed to {target_status_name}"
        if reason:
            log_title += f". Reason: {reason}"
            
        l = WorkOrderLog(work_order_id=wo.id, user_id=current_user.id, log_text=log_title)
        db.session.add(l)
            
        db.session.commit()
        flash(f'Work order status changed to {target_status_name}.', 'success')
        
    return redirect(url_for('work_orders.edit', id=wo.id))

@work_orders_bp.route('/<int:id>/add_part', methods=['POST'])
@login_required
def add_part(id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        return redirect(url_for('work_orders.index'))
        
    if wo.current_status and wo.current_status.name in ['Solved', 'Completed', 'Closed', 'Incompleted']:
        flash('Cannot add parts to a Work Order that is already solved or completed.', 'warning')
        return redirect(url_for('work_orders.edit', id=wo.id))
        
    part_id = request.form.get('part_id')
    qty = float(request.form.get('qty', 1))
    
    if not part_id or qty <= 0:
        flash('Invalid part or quantity.', 'danger')
        return redirect(url_for('work_orders.edit', id=wo.id))
        
    # Verify stock availability at the site
    stock = StockLevel.query.filter_by(part_id=part_id, site_id=wo.asset.site_id).first()
    if not stock or stock.qty_on_hand < qty:
        flash('Insufficient stock available at this site.', 'danger')
        return redirect(url_for('work_orders.edit', id=wo.id))
        
    # Deduct stock
    qty_before = stock.qty_on_hand
    stock.qty_on_hand -= qty
    qty_after = stock.qty_on_hand
    
    # Log Transaction
    from models import StockTransaction
    from datetime import datetime
    tx_code = f"TRX-{datetime.now().strftime('%Y%m%d%H%M%S')}"
    
    transaction = StockTransaction(
        transaction_code=tx_code,
        part_id=part_id,
        site_id=wo.asset.site_id,
        user_id=current_user.id,
        action='WORK ORDER',
        qty_before=qty_before,
        qty_after=qty_after,
        qty_change=-qty,
        notes=f"Used in WO: {wo.ticket_code}"
    )
    db.session.add(transaction)
    
    # Check if part already exists in this WO
    wo_part = WorkOrderPart.query.filter_by(work_order_id=wo.id, part_id=part_id).first()
    if wo_part:
        wo_part.quantity_used += qty
    else:
        wo_part = WorkOrderPart(work_order_id=wo.id, part_id=part_id, quantity_used=qty)
        db.session.add(wo_part)
        
    db.session.commit()
    flash(f'{qty} units of part added successfully. Stock deducted.', 'success')
    return redirect(url_for('work_orders.edit', id=wo.id))

@work_orders_bp.route('/<int:id>/remove_part/<int:part_id>', methods=['POST'])
@login_required
def remove_part(id, part_id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        return redirect(url_for('work_orders.index'))
        
    wo_part = WorkOrderPart.query.filter_by(work_order_id=wo.id, part_id=part_id).first_or_404()
    
    if wo.current_status and wo.current_status.name in ['Solved', 'Completed', 'Closed', 'Incompleted']:
        flash('Cannot remove parts from a Work Order that is already solved or completed.', 'warning')
        return redirect(url_for('work_orders.edit', id=wo.id))
    
    # Restore stock
    stock = StockLevel.query.filter_by(part_id=part_id, site_id=wo.asset.site_id).first()
    qty_before = 0
    if stock:
        qty_before = stock.qty_on_hand
        stock.qty_on_hand += wo_part.quantity_used
    else:
        # If stock level entry was deleted somehow, recreate it
        stock = StockLevel(part_id=part_id, site_id=wo.asset.site_id, qty_on_hand=wo_part.quantity_used)
        db.session.add(stock)
        
    qty_after = stock.qty_on_hand
    
    # Log Transaction
    from models import StockTransaction
    from datetime import datetime
    tx_code = f"TRX-{datetime.now().strftime('%Y%m%d%H%M%S')}"
    
    transaction = StockTransaction(
        transaction_code=tx_code,
        part_id=part_id,
        site_id=wo.asset.site_id,
        user_id=current_user.id,
        action='RETURN',
        qty_before=qty_before,
        qty_after=qty_after,
        qty_change=wo_part.quantity_used,
        notes=f"Returned from WO: {wo.ticket_code}"
    )
    db.session.add(transaction)
        
    db.session.delete(wo_part)
    db.session.commit()
    
    flash('Part removed from Work Order. Stock refunded.', 'info')
    return redirect(url_for('work_orders.edit', id=wo.id))

@work_orders_bp.route('/<int:id>/delete_attachment/<int:attachment_id>', methods=['POST'])
@login_required
def delete_attachment(id, attachment_id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        return redirect(url_for('work_orders.index'))
        
    from models import WorkOrderAttachment
    att = WorkOrderAttachment.query.filter_by(id=attachment_id, work_order_id=wo.id).first_or_404()
    
    # Optional: Delete actual file from disk here if desired
    # file_path_sys = os.path.join(current_app.root_path, 'static', att.file_path)
    # if os.path.exists(file_path_sys):
    #     os.remove(file_path_sys)
        
    db.session.delete(att)
    db.session.commit()
    
    flash('Attachment removed successfully.', 'info')
    return redirect(url_for('work_orders.edit', id=wo.id))

@work_orders_bp.route('/<int:id>/add_log', methods=['POST'])
@login_required
def add_log(id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        return redirect(url_for('work_orders.index'))
        
    log_text = request.form.get('log_text')
    
    if not log_text or log_text.strip() == '':
        flash('Log cannot be empty.', 'danger')
        return redirect(url_for('work_orders.edit', id=wo.id))
        
    new_log = WorkOrderLog(
        work_order_id=wo.id,
        user_id=current_user.id,
        log_text=log_text.strip()
    )
    db.session.add(new_log)
    db.session.commit()
    
    flash('Work log added successfully.', 'success')
    return redirect(url_for('work_orders.edit', id=wo.id) + '#log')

@work_orders_bp.route('/<int:id>/delete', methods=['POST'])
@login_required
@supervisor_or_admin_required
def delete(id):
    wo = WorkOrder.query.get_or_404(id)
    db.session.delete(wo)
    db.session.commit()
    flash('Work Order deleted successfully!', 'success')
    return redirect(url_for('work_orders.index'))

@work_orders_bp.route('/<int:id>/procedure/<int:proc_id>/complete', methods=['POST'])
@login_required
def complete_procedure(id, proc_id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        return redirect(url_for('work_orders.index'))
        
    proc = WorkOrderProcedure.query.filter_by(id=proc_id, work_order_id=wo.id).first_or_404()
    
    if not wo.current_status or wo.current_status.control_type != 'Active':
        flash('Work Order must be In Progress (Active) to modify tasks.', 'danger')
        return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')

    action = request.form.get('action', 'complete')
    if action == 'uncomplete':
        proc.is_completed = False
        proc.completed_at = None
        proc.completed_by_id = None
        db.session.commit()
        return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')

    if action == 'append_photo':
        files = request.files.getlist('proc_attachment')
        if not files or all(f.filename == '' for f in files):
            flash(f"No file selected to append.", 'danger')
            return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')
            
        saved_paths = []
        import uuid
        for file in files:
            if file and file.filename != '':
                filename = secure_filename(f"wo{wo.id}_proc{proc.id}_{uuid.uuid4().hex[:6]}_{file.filename}")
                uploads_dir = os.path.join(current_app.root_path, 'static', 'uploads')
                os.makedirs(uploads_dir, exist_ok=True)
                file_path_sys = os.path.join(uploads_dir, filename)
                file.save(file_path_sys)
                saved_paths.append(f'uploads/{filename}')
        
        if saved_paths:
            new_paths = '|'.join(saved_paths)
            if proc.attachment_path:
                proc.attachment_path += '|' + new_paths
            else:
                proc.attachment_path = new_paths
            db.session.commit()
            flash(f"Photos appended successfully.", 'success')
            
        return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')

    # Handle Skip / Fail
    if action == 'fail':
        notes = request.form.get('notes')
        if not notes or not notes.strip():
            flash(f"Note is required when skipping procedure '{proc.name}'.", 'danger')
            return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')
        
        proc.status = 'Failed'
        proc.notes = notes
        proc.is_completed = True
        proc.completed_at = datetime.now(timezone.utc).replace(tzinfo=None)
        proc.completed_by_id = current_user.id
        db.session.commit()
        flash(f"Procedure '{proc.name}' was skipped/failed successfully with a note.", 'warning')
        return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')

    if proc.requires_attachment:
        if 'proc_attachment' not in request.files:
            flash(f"Procedure '{proc.name}' requires a photo attachment.", 'danger')
            return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')
            
        files = request.files.getlist('proc_attachment')
        if not files or all(f.filename == '' for f in files):
            flash(f"No file selected for procedure '{proc.name}'.", 'danger')
            return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')
            
        saved_paths = []
        for file in files:
            if file and file.filename != '':
                filename = secure_filename(f"wo{wo.id}_proc{proc.id}_{file.filename}")
                uploads_dir = os.path.join(current_app.root_path, 'static', 'uploads')
                os.makedirs(uploads_dir, exist_ok=True)
                file_path_sys = os.path.join(uploads_dir, filename)
                file.save(file_path_sys)
                saved_paths.append(f'uploads/{filename}')
        
        if saved_paths:
            proc.attachment_path = '|'.join(saved_paths)
        
    proc.is_completed = True
    proc.status = 'Done'
    proc.notes = None
    proc.completed_at = datetime.now(timezone.utc).replace(tzinfo=None)
    proc.completed_by_id = current_user.id
    db.session.commit()
    flash(f"Procedure '{proc.name}' completed successfully.", 'success')
    return redirect(url_for('work_orders.edit', id=wo.id) + '#procedures')

@work_orders_bp.route('/<int:id>/checklist_param/<int:chk_id>/save', methods=['POST'])
@login_required
def save_checklist_parameter(id, chk_id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        return redirect(url_for('work_orders.index'))
        
    if not wo.current_status or wo.current_status.control_type != 'Active':
        flash('Work Order must be In Progress (Active) to modify checklists.', 'danger')
        return redirect(url_for('work_orders.edit', id=wo.id) + '#checklist_panel')
        
    chk = WorkOrderChecklistParameter.query.filter_by(id=chk_id, work_order_id=wo.id).first_or_404()
    
    val = request.form.get('value')
    if val:
        chk.value = val
    chk.note = request.form.get('note')
    
    db.session.commit()
    flash(f"Checklist parameter '{chk.parameter}' saved.", 'success')
    return redirect(url_for('work_orders.edit', id=wo.id) + '#checklist_panel')

@work_orders_bp.route('/<int:id>/checklist/bulk_save', methods=['POST'])
@login_required
def bulk_save_checklist(id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        return redirect(url_for('work_orders.index'))
        
    if not wo.current_status or wo.current_status.control_type != 'Active':
        flash('Work Order must be In Progress (Active) to modify checklists.', 'danger')
        return redirect(url_for('work_orders.edit', id=wo.id) + '#checklist_panel')
        
    for chk in wo.checklist_parameters:
        val = request.form.get(f'value_{chk.id}')
        note = request.form.get(f'note_{chk.id}')
        # We update even if empty if it was sent in the form
        if val is not None:
            chk.value = val
        if note is not None:
            chk.note = note
            
    db.session.commit()
    flash('All checklist parameters saved successfully.', 'success')
    return redirect(url_for('work_orders.edit', id=wo.id) + '#checklist_panel')

@work_orders_bp.route('/<int:id>/procedure/add', methods=['POST'])
@login_required
def add_wo_procedure(id):
    if current_user.role not in ['Admin', 'Supervisor']:
        flash('Permission denied.', 'danger')
        return redirect(url_for('work_orders.edit', id=id))
        
    wo = WorkOrder.query.get_or_404(id)
    name = request.form.get('name')
    try:
        estimated_minutes = float(request.form.get('estimated_minutes', 0))
    except:
        estimated_minutes = 0
    requires_attachment = 'requires_attachment' in request.form
    
    if name:
        new_proc = WorkOrderProcedure(work_order=wo, tasklist_name='Custom Task', name=name, requires_attachment=requires_attachment, estimated_minutes=estimated_minutes)
        db.session.add(new_proc)
        db.session.commit()
        flash(f'Custom procedure "{name}" added successfully.', 'success')
        
    return redirect(url_for('work_orders.edit', id=id) + '#procedures_panel')

@work_orders_bp.route('/<int:id>/procedure/<int:proc_id>/edit', methods=['POST'])
@login_required
def edit_wo_procedure(id, proc_id):
    if current_user.role not in ['Admin', 'Supervisor']:
        flash('Permission denied.', 'danger')
        return redirect(url_for('work_orders.edit', id=id))
        
    wo = WorkOrder.query.get_or_404(id)
    proc = WorkOrderProcedure.query.filter_by(id=proc_id, work_order_id=wo.id).first_or_404()
    
    name = request.form.get('name')
    if name:
        proc.name = name
    try:
        proc.estimated_minutes = float(request.form.get('estimated_minutes', 0))
    except:
        pass
    proc.requires_attachment = 'requires_attachment' in request.form
    
    db.session.commit()
    flash('Procedure updated successfully.', 'success')
    return redirect(url_for('work_orders.edit', id=id) + '#procedures_panel')

@work_orders_bp.route('/<int:id>/procedure/<int:proc_id>/delete', methods=['POST'])
@login_required
def delete_wo_procedure(id, proc_id):
    if current_user.role not in ['Admin', 'Supervisor']:
        flash('Permission denied.', 'danger')
        return redirect(url_for('work_orders.edit', id=id))
        
    wo = WorkOrder.query.get_or_404(id)
    proc = WorkOrderProcedure.query.filter_by(id=proc_id, work_order_id=wo.id).first_or_404()
    
    db.session.delete(proc)
    db.session.commit()
    flash('Procedure deleted successfully.', 'success')
    return redirect(url_for('work_orders.edit', id=id) + '#procedures_panel')

@work_orders_bp.route('/<int:id>/checklist/add', methods=['POST'])
@login_required
def add_wo_checklist(id):
    if current_user.role not in ['Admin', 'Supervisor']:
        flash('Permission denied.', 'danger')
        return redirect(url_for('work_orders.edit', id=id))
        
    wo = WorkOrder.query.get_or_404(id)
    parameter = request.form.get('parameter')
    standard = request.form.get('standard')
    
    if parameter and standard:
        new_chk = WorkOrderChecklistParameter(work_order=wo, checklist_name='Custom Check', parameter=parameter, standard=standard)
        db.session.add(new_chk)
        db.session.commit()
        flash('Custom parameter added successfully.', 'success')
        
    return redirect(url_for('work_orders.edit', id=id) + '#checklist_panel')

@work_orders_bp.route('/<int:id>/checklist/<int:chk_id>/edit', methods=['POST'])
@login_required
def edit_wo_checklist(id, chk_id):
    if current_user.role not in ['Admin', 'Supervisor']:
        flash('Permission denied.', 'danger')
        return redirect(url_for('work_orders.edit', id=id))
        
    wo = WorkOrder.query.get_or_404(id)
    chk = WorkOrderChecklistParameter.query.filter_by(id=chk_id, work_order_id=wo.id).first_or_404()
    
    parameter = request.form.get('parameter')
    standard = request.form.get('standard')
    
    if parameter:
        chk.parameter = parameter
    if standard:
        chk.standard = standard
        
    db.session.commit()
    flash('Checklist parameter updated successfully.', 'success')
    return redirect(url_for('work_orders.edit', id=id) + '#checklist_panel')

@work_orders_bp.route('/<int:id>/checklist/<int:chk_id>/delete', methods=['POST'])
@login_required
def delete_wo_checklist(id, chk_id):
    if current_user.role not in ['Admin', 'Supervisor']:
        flash('Permission denied.', 'danger')
        return redirect(url_for('work_orders.edit', id=id))
        
    wo = WorkOrder.query.get_or_404(id)
    chk = WorkOrderChecklistParameter.query.filter_by(id=chk_id, work_order_id=wo.id).first_or_404()
    
    db.session.delete(chk)
    db.session.commit()
    flash('Checklist parameter deleted successfully.', 'success')
    return redirect(url_for('work_orders.edit', id=id) + '#checklist_panel')
@login_required
def export_pptx(id):
    wo = WorkOrder.query.get_or_404(id)
    if not check_wo_access(wo):
        flash('Access denied.', 'danger')
        return redirect(url_for('work_orders.index'))
        
    prs = Presentation()
    prs.slide_width = Cm(21.0)
    prs.slide_height = Cm(29.7)
    blank_layout = prs.slide_layouts[6]
    
    # --- Helper to format cells ---
    def format_cell(cell, text, bold=False, size=8, align=PP_ALIGN.LEFT, fill_color=None):
        cell.text = str(text) if text is not None else ""
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.name = "Helvetica"
            p.alignment = align
        if fill_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

    def fmt_dt(dt, convert=True):
        if dt:
            if convert:
                dt_wib = dt + timedelta(hours=7)
                return dt_wib.strftime('%d-%m-%Y %H:%M')
            return dt.strftime('%d-%m-%Y %H:%M')
        return '-'

    bg_color = RGBColor(220, 220, 220)
    
    # --- 0. COVER PAGE ---
    slide = prs.slides.add_slide(blank_layout)
    title_image_path = os.path.join(current_app.root_path, 'static', 'images', 'judul.png')
    if os.path.exists(title_image_path):
        slide.shapes.add_picture(title_image_path, Cm(0), Cm(0), width=prs.slide_width, height=prs.slide_height)
        
    site_name = wo.asset.site.name if wo.asset and getattr(wo.asset, 'site', None) else "Unknown Site"
    wo_date = wo.suggested_start_date or wo.start_date or datetime.now()
    month_map = {
        1: "JANUARI", 2: "FEBRUARI", 3: "MARET", 4: "APRIL", 5: "MEI", 6: "JUNI",
        7: "JULI", 8: "AGUSTUS", 9: "SEPTEMBER", 10: "OKTOBER", 11: "NOVEMBER", 12: "DESEMBER"
    }
    periode_str = f"{month_map.get(wo_date.month, '')} {wo_date.year}"
    
    mt_type = wo.maintenance_type.upper() if wo.maintenance_type else ""
    if "MAINTENANCE" not in mt_type:
        mt_type += " MAINTENANCE"
        
    line1 = "LAPORAN WORK ORDER"
    line1_sub = mt_type
    line2 = site_name.upper()
    line3 = wo.asset.name.upper() if wo.asset else ''
    line4 = wo.code
    
    txBox = slide.shapes.add_textbox(Cm(2), Cm(18), Cm(17), Cm(8))
    tf = txBox.text_frame
    
    p = tf.paragraphs[0]
    p.text = line1
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.name = "Helvetica"
    
    p = tf.add_paragraph()
    p.text = line1_sub
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.name = "Helvetica"
    
    tf.add_paragraph() # space
    
    p = tf.add_paragraph()
    p.text = line2
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.name = "Helvetica"
    
    p = tf.add_paragraph()
    p.text = line3
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.name = "Helvetica"
    
    p = tf.add_paragraph()
    p.text = line4
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.name = "Helvetica"
    
    # --- END COVER PAGE ---
    
    # --- CONTENT PAGES ---
    def add_content_slide():
        s = prs.slides.add_slide(blank_layout)
        
        # 1. NEW HEADER
        header_shape = s.shapes.add_table(4, 4, Cm(1), Cm(1), Cm(19), Cm(2.4))
        ht = header_shape.table
        ht.columns[0].width = Cm(3.5)
        ht.columns[1].width = Cm(10.5)
        ht.columns[2].width = Cm(2.0)
        ht.columns[3].width = Cm(3.0)
        
        ht.cell(0,0).merge(ht.cell(3,0))
        ht.cell(0,1).merge(ht.cell(3,1))
        
        logo_path = os.path.join(current_app.root_path, 'static', 'images', 'Logo Jaya Teknik.png')
        if os.path.exists(logo_path):
            s.shapes.add_picture(logo_path, Cm(1.2), Cm(1.2), width=Cm(3.1))
            
        format_cell(ht.cell(0,1), "REPORT WORK ORDER", bold=True, size=14, align=PP_ALIGN.CENTER)
        
        format_cell(ht.cell(0,2), "No.Dok")
        format_cell(ht.cell(0,3), "RP - JT - 26")
        format_cell(ht.cell(1,2), "Ref.")
        format_cell(ht.cell(1,3), "7; 7.1.3 & 7.1.4")
        format_cell(ht.cell(2,2), "Rev.")
        format_cell(ht.cell(2,3), "Original")
        format_cell(ht.cell(3,2), "Tanggal")
        format_cell(ht.cell(3,3), datetime.now().strftime("%d %b %Y"))
        
        return s
        
    s = add_content_slide()
    current_y = Cm(3.6)
    
    # Details Body
    form_val = getattr(wo, 'maintenance_type', 'SERVICE')
    nomor_val = wo.code or '-'
    asset_val = f"{wo.asset.name} ({wo.asset.code})" if wo.asset else '-'
    resp_val = ", ".join([a.name for a in wo.assignees]) if getattr(wo, 'assignees', None) else "N/A"
    
    proj_val = (wo.asset.project_code if wo.asset and wo.asset.project_code else wo.project_code) or '-'
    loc_val = '-'
    if wo.asset and hasattr(wo.asset, 'location') and wo.asset.location:
        loc_val = wo.asset.location.name if hasattr(wo.asset.location, 'name') else str(wo.asset.location)
        
    sched_val = f"{fmt_dt(wo.suggested_start_date)} s/d {fmt_dt(wo.suggested_completion_date)}"
    act_val = f"{fmt_dt(wo.start_date)} s/d {fmt_dt(wo.end_date)}"
    
    dt_shape = s.shapes.add_table(4, 6, Cm(1), current_y, Cm(19), Cm(2.4))
    dt = dt_shape.table
    dt.columns[0].width = Cm(2.5)
    dt.columns[1].width = Cm(0.5)
    dt.columns[2].width = Cm(6.5)
    dt.columns[3].width = Cm(2.5)
    dt.columns[4].width = Cm(0.5)
    dt.columns[5].width = Cm(6.5)
    
    data_rows = [
        ["Form", ":", form_val, "Project", ":", proj_val],
        ["Nomor", ":", nomor_val, "Location", ":", loc_val],
        ["Asset", ":", asset_val, "Schedule", ":", sched_val],
        ["Responsible", ":", resp_val, "Actual Date", ":", act_val]
    ]
    for r_idx, row_data in enumerate(data_rows):
        for c_idx, val in enumerate(row_data):
            format_cell(dt.cell(r_idx, c_idx), val, size=9)
            
    current_y += Cm(2.6)
    
    # 1.5 TECHNICIAN
    tech_shape = s.shapes.add_table(3, 1, Cm(1), current_y, Cm(19), Cm(1.8))
    tech = tech_shape.table
    format_cell(tech.cell(0,0), "TECHNICIAN", bold=True, size=10, align=PP_ALIGN.CENTER, fill_color=bg_color)
    format_cell(tech.cell(1,0), "Name", bold=True, size=10, align=PP_ALIGN.CENTER, fill_color=bg_color)
    format_cell(tech.cell(2,0), resp_val, size=9)
    
    current_y += Cm(2.2)
    
    # 2. CHECKING REPORT
    chk_params = wo.checklist_parameters
    chk_rows = len(chk_params) + 2
    
    # If the table is too long, we might need a new slide. For simplicity, we just put it on the current slide.
    chk_shape = s.shapes.add_table(chk_rows, 5, Cm(1), current_y, Cm(19), Cm(0.6 * chk_rows))
    chk = chk_shape.table
    chk.columns[0].width = Cm(7.5)
    chk.columns[1].width = Cm(2.5)
    chk.columns[2].width = Cm(2.5)
    chk.columns[3].width = Cm(2.5)
    chk.columns[4].width = Cm(4.0)
    
    chk.cell(0,0).merge(chk.cell(0,4))
    format_cell(chk.cell(0,0), "CHECKING REPORT", bold=True, size=10, align=PP_ALIGN.CENTER, fill_color=bg_color)
    
    headers = ["Description (Unit Check)", "Standard", "Actual", "Check", "Note"]
    for i, h in enumerate(headers):
        format_cell(chk.cell(1,i), h, bold=True, size=8, align=PP_ALIGN.CENTER, fill_color=bg_color)
        
    for r_idx, param in enumerate(chk_params):
        desc = (param.parameter[:40] + '...') if len(param.parameter) > 42 else param.parameter
        actual_val = str(param.value) if param.value else "-"
        check_val = "v" if actual_val != "-" else "-"
        note_val = str(param.note) if param.note else ""
        
        format_cell(chk.cell(r_idx+2, 0), desc)
        format_cell(chk.cell(r_idx+2, 1), str(param.standard)[:15], align=PP_ALIGN.CENTER)
        format_cell(chk.cell(r_idx+2, 2), actual_val[:15], align=PP_ALIGN.CENTER)
        format_cell(chk.cell(r_idx+2, 3), "OK" if check_val == "v" else "-", align=PP_ALIGN.CENTER)
        format_cell(chk.cell(r_idx+2, 4), note_val[:20])
        
    current_y += Cm(0.6 * chk_rows) + Cm(0.4)
    
    # Check if we need a new slide for Tasklist
    if current_y > Cm(22):
        s = add_content_slide()
        current_y = Cm(3.6)
        
    # 3. TASKLIST REPORT
    procs = wo.procedures
    proc_rows = len(procs) + 2
    
    proc_shape = s.shapes.add_table(proc_rows, 3, Cm(1), current_y, Cm(19), Cm(0.6 * proc_rows))
    proc_tbl = proc_shape.table
    proc_tbl.columns[0].width = Cm(7.0)
    proc_tbl.columns[1].width = Cm(9.0)
    proc_tbl.columns[2].width = Cm(3.0)
    
    proc_tbl.cell(0,0).merge(proc_tbl.cell(0,2))
    format_cell(proc_tbl.cell(0,0), "TASKLIST REPORT", bold=True, size=10, align=PP_ALIGN.CENTER, fill_color=bg_color)
    
    headers = ["Step Name", "Description", "State"]
    for i, h in enumerate(headers):
        format_cell(proc_tbl.cell(1,i), h, bold=True, size=8, align=PP_ALIGN.CENTER, fill_color=bg_color)
        
    for r_idx, proc in enumerate(procs):
        step_name = proc.name or ''
        desc = getattr(proc, 'description', '') or ''
        state = "Done" if getattr(proc, 'is_completed', False) else "Pending"
        
        format_cell(proc_tbl.cell(r_idx+2, 0), step_name)
        format_cell(proc_tbl.cell(r_idx+2, 1), desc)
        format_cell(proc_tbl.cell(r_idx+2, 2), state, align=PP_ALIGN.CENTER)
        
    current_y += Cm(0.6 * proc_rows) + Cm(0.4)
    
    # 3.6 SIGNATURES
    if current_y > Cm(24):
        s = add_content_slide()
        current_y = Cm(3.6)
        
    first_tech = wo.assignees[0] if getattr(wo, 'assignees', []) else None
    tech_name = first_tech.name if first_tech else "....."
    tech_role = first_tech.role if first_tech else "....."
    site_name = wo.asset.site.name.upper() if wo.asset and getattr(wo.asset, 'site', None) else "....."
    c_name_print = getattr(wo, 'customer_name', None) or "....."
    c_title_print = getattr(wo, 'customer_title', None) or "....."
    
    sig_shape = s.shapes.add_table(3, 3, Cm(1.5), current_y, Cm(18), Cm(4))
    sig = sig_shape.table
    sig.columns[0].width = Cm(7.0)
    sig.columns[1].width = Cm(4.0)
    sig.columns[2].width = Cm(7.0)
    
    format_cell(sig.cell(0,0), "PIHAK PERTAMA,\n" + site_name, bold=True, size=9)
    format_cell(sig.cell(0,2), "PIHAK KEDUA,\nPT JAYA TEKNIK INDONESIA", bold=True, size=9)
    
    # We leave cell(1,x) for signature images
    
    format_cell(sig.cell(2,0), f"Nama : {c_name_print}\nJabatan : {c_title_print}", size=9)
    format_cell(sig.cell(2,2), f"Nama : {tech_name}\nJabatan : {tech_role}", size=9)
    
    # Add Digital Signature QR if exists
    from models import DigitalSignature
    ds = DigitalSignature.query.filter_by(work_order_id=wo.id).first()
    
    def save_temp_image(b64_str):
        if not b64_str or ',' not in b64_str: return None
        try:
            _, encoded = b64_str.split(",", 1)
            data = base64.b64decode(encoded)
            fd, temp_path = tempfile.mkstemp(suffix='.png')
            with os.fdopen(fd, 'wb') as f:
                f.write(data)
            return temp_path
        except:
            return None
            
    if ds:
        import qrcode
        qr = qrcode.QRCode(version=1, box_size=10, border=1)
        verify_url = f"{request.url_root.rstrip('/')}/verify/doc?id={ds.id}"
        qr.add_data(verify_url)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        
        fd, temp_path = tempfile.mkstemp(suffix='.png')
        with os.fdopen(fd, 'wb') as f:
            img.save(f)
            
        s.shapes.add_picture(temp_path, Cm(9), current_y + Cm(1), height=Cm(1.8))
        os.remove(temp_path)
        
        format_cell(sig.cell(2,1), "DIGITALLY SIGNED\n& VERIFIED", bold=True, size=6, align=PP_ALIGN.CENTER)
    else:
        # Add normal signatures
        c_img_path = save_temp_image(getattr(wo, 'customer_signature', None))
        if c_img_path:
            try:
                s.shapes.add_picture(c_img_path, Cm(1.5), current_y + Cm(1.2), height=Cm(1.6))
                os.remove(c_img_path)
            except: pass
            
        t_img_path = save_temp_image(getattr(wo, 'technician_signature', None))
        if t_img_path:
            try:
                s.shapes.add_picture(t_img_path, Cm(12.5), current_y + Cm(1.2), height=Cm(1.6))
                os.remove(t_img_path)
            except: pass

    current_y += Cm(4.5)
    
    # 4. PHOTOS
    images = []
    for proc in wo.procedures:
        if proc.attachment_path:
            for path in proc.attachment_path.split('|'):
                if path:
                    images.append({"caption": f"Lampiran {proc.name[:50]}", "path": path})
    for att in wo.attachments:
        images.append({"caption": f"Lampiran {att.file_name}", "path": att.file_path})
        
    if images:
        s = add_content_slide()
        current_y = Cm(3.6)
        
        title_shape = s.shapes.add_table(1, 1, Cm(1), current_y, Cm(19), Cm(0.6))
        format_cell(title_shape.table.cell(0,0), "LAMPIRAN FOTO", bold=True, size=10, align=PP_ALIGN.CENTER, fill_color=bg_color)
        current_y += Cm(1.0)
        
        box_width = Cm(6.33)
        box_height = Cm(9.5)
        
        for i in range(0, len(images), 3):
            if current_y + box_height > Cm(28):
                s = add_content_slide()
                current_y = Cm(3.6)
                
            row_images = images[i:i+3]
            for j, img_dict in enumerate(row_images):
                x_pos = Cm(1) + j * box_width
                
                # Draw caption
                tx = s.shapes.add_textbox(x_pos, current_y, box_width, Cm(1))
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = img_dict['caption']
                p.font.size = Pt(7)
                p.alignment = PP_ALIGN.CENTER
                
                # Draw image
                full_path = os.path.join(current_app.root_path, 'static', img_dict['path'])
                if os.path.exists(full_path):
                    try:
                        from PIL import Image, ImageOps
                        with Image.open(full_path) as im:
                            im = ImageOps.exif_transpose(im)
                            w, h = im.size
                            if w > h:
                                im = im.rotate(270, expand=True)
                                w, h = im.size
                                
                            target_aspect = 3/4
                            aspect = w/h
                            if aspect > target_aspect:
                                new_w = int(h * target_aspect)
                                left = (w - new_w) / 2
                                im = im.crop((left, 0, left + new_w, h))
                            else:
                                new_h = int(w / target_aspect)
                                top = (h - new_h) / 2
                                im = im.crop((0, top, w, top + new_h))
                                
                            im = im.resize((300, 400))
                            
                            tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
                            if im.mode != 'RGB':
                                im = im.convert('RGB')
                            im.save(tmp.name, 'JPEG')
                            tmp.close()
                            
                            s.shapes.add_picture(tmp.name, x_pos + Cm(0.3), current_y + Cm(1.3), width=Cm(5.7))
                            os.unlink(tmp.name)
                    except Exception as e:
                        print(f"Error processing image {full_path}: {e}")
                        
            current_y += box_height
            
    # Output to user
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    filename = f"WO_{wo.code}_{datetime.now().strftime('%Y%m%d')}.pptx"
    return send_file(pptx_io, as_attachment=True, download_name=filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
